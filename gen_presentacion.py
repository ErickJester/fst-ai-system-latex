# -*- coding: utf-8 -*-
"""
Genera la presentacion ejecutiva de defensa de tesis TT 2026-B066.
Usa diagramas de la carpeta diagramas/.
Estructura: 22 slides segun el TOC solicitado.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ==========================================================
# Paleta de colores (Ocean Academic)
# ==========================================================
NAVY       = RGBColor(0x0B, 0x24, 0x47)
MIDNIGHT   = RGBColor(0x19, 0x37, 0x6D)
BLUE       = RGBColor(0x57, 0x6C, 0xBC)
ICE        = RGBColor(0xCA, 0xDC, 0xFC)
BG_LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x00, 0xC8, 0x96)
WARN       = RGBColor(0xE9, 0x6B, 0x37)
TEXT_DARK  = RGBColor(0x1A, 0x20, 0x2C)
TEXT_MUTE  = RGBColor(0x64, 0x74, 0x8B)
GRID       = RGBColor(0xE2, 0xE8, 0xF0)

BASE = r"C:\Users\ID634\Desktop\proyectos\latexFST\fst-ai-system-latex"
DIAG = os.path.join(BASE, "diagramas")
LOGOS = os.path.join(BASE, "logos")
OUT = os.path.join(BASE, "Presentacion_TT_2026-B066.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = 13.333
SH = 7.5
BLANK = prs.slide_layouts[6]

# ==========================================================
# Helpers
# ==========================================================
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_noline(shape):
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, color, line=False, line_color=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(s, color)
    if line and line_color is not None:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    else:
        set_noline(s)
    s.shadow.inherit = False
    return s

def add_rounded(slide, x, y, w, h, color, line=False, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(s, color)
    if line and line_color is not None:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.75)
    else:
        set_noline(s)
    return s

def add_oval(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(s, color)
    set_noline(s)
    return s

def add_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, size=14, color=TEXT_DARK, bullet_color=None,
                font="Calibri", line_sp=1.15, marker="•"):
    if bullet_color is None:
        bullet_color = ACCENT
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    first = True
    for it in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_sp
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"{marker}  "
        r1.font.name = font
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = it
        r2.font.name = font
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb

def add_line(slide, x1, y1, x2, y2, color, weight=1.5):
    l = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    l.line.color.rgb = color
    l.line.width = Pt(weight)
    return l

def slide_header_bar(slide, title, subtitle=None, chapter_tag=None):
    add_rect(slide, 0, 0, 0.28, SH, NAVY)
    add_text(slide, 0.55, 0.35, 11.5, 0.7, title, size=28, bold=True, color=NAVY, font="Georgia")
    if subtitle is not None:
        add_text(slide, 0.55, 0.95, 11.5, 0.4, subtitle, size=13, color=TEXT_MUTE, italic=True)
    if chapter_tag is not None:
        add_rounded(slide, 11.4, 0.4, 1.6, 0.35, ACCENT)
        add_text(slide, 11.4, 0.4, 1.6, 0.35, chapter_tag, size=10, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_line(slide, 0.55, 1.45, 4.2, 1.45, ACCENT, weight=2.2)

def footer_bar(slide, idx, total):
    add_rect(slide, 0, SH-0.32, SW, 0.32, NAVY)
    add_text(slide, 0.55, SH-0.32, 8, 0.32,
             "TT 2026-B066  ·  ESCOM-IPN  ·  Sistema FST con IA",
             size=9, color=ICE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW-1.5, SH-0.32, 1, 0.32, f"{idx} / {total}",
             size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


# ==========================================================
# SLIDE 1 — Portada
# ==========================================================
def slide_portada():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, 0, SW, 0.15, ACCENT)
    add_rect(s, 0, SH-0.15, SW, 0.15, ACCENT)

    if os.path.exists(os.path.join(LOGOS, "ipn_logo.png")):
        s.shapes.add_picture(os.path.join(LOGOS, "ipn_logo.png"),
                             Inches(0.55), Inches(0.4), height=Inches(0.9))
    if os.path.exists(os.path.join(LOGOS, "logo_escom.png")):
        s.shapes.add_picture(os.path.join(LOGOS, "logo_escom.png"),
                             Inches(SW-1.55), Inches(0.4), height=Inches(0.9))

    add_text(s, 0, 1.6, SW, 0.4,
             "INSTITUTO POLITÉCNICO NACIONAL  ·  ESCUELA SUPERIOR DE CÓMPUTO",
             size=12, bold=True, color=ICE, align=PP_ALIGN.CENTER)
    add_text(s, 0, 2.0, SW, 0.3, "Trabajo Terminal 2026-B066",
             size=11, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

    add_text(s, 0.8, 2.6, SW-1.6, 1.5,
             "Prototipo de un sistema web\nde análisis conductual asistido por IA",
             size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Georgia")

    add_text(s, 0.8, 3.95, SW-1.6, 0.6,
             "Modelo de conducta depresiva en rata mediante nado forzado",
             size=17, color=ICE, align=PP_ALIGN.CENTER, italic=True)

    add_line(s, 4.5, 4.8, 8.85, 4.8, ACCENT, weight=2)

    add_text(s, 0, 4.95, SW, 0.35, "Presentan",
             size=11, color=ICE, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, 0, 5.25, SW, 0.4, "Frausto Robles Ángel Ali    ·    Rodríguez Verdín Sandoval Vanesa",
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, 0, 5.85, SW, 0.3, "Directores",
             size=10, color=ICE, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, 0, 6.1, SW, 0.35, "Dra. Cordero López Martha Rosa    ·    Dr. Salas Ramírez Israel",
             size=12, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, 0, 6.7, SW, 0.35, "Mayo 2026",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font="Georgia")

# ==========================================================
# SLIDE 2 — Agenda (6 secciones)
# ==========================================================
def slide_agenda(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Agenda", "Contenido de la presentación")

    items = [
        ("1", "Introducción y problemática",
         "Contexto FST · problema · propuesta de solución"),
        ("2", "Objetivos y justificación",
         "Metas técnicas · ISO/IEC 25010 · alcance del prototipo"),
        ("3", "Fundamentos",
         "Estado del arte · marcos teóricos (conductual, visión, DL)"),
        ("4", "Metodología e ingeniería de requisitos",
         "Scrum · RF · RNF · reglas de negocio"),
        ("5", "Diseño del sistema",
         "Riesgos · arquitectura · pipeline · BD · vistas"),
        ("6", "Validación y cierre",
         "Métricas · conclusiones · trabajo futuro · bibliografía"),
    ]
    top = 1.9
    row_h = 0.78
    for i, (n, t, d) in enumerate(items):
        y = top + i*row_h
        add_oval(s, 0.8, y, 0.58, 0.58, NAVY)
        add_text(s, 0.8, y, 0.58, 0.58, n, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, 1.6, y-0.02, 8, 0.4, t, size=17, bold=True, color=NAVY)
        add_text(s, 1.6, y+0.34, 11, 0.35, d, size=12, color=TEXT_MUTE, italic=True)

    footer_bar(s, 2, total)

# ==========================================================
# SLIDE 3 — Introducción (antecedentes FST + neurociencia preclínica)
# ==========================================================
def slide_introduccion(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Introducción",
                     "Antecedentes del nado forzado y neurociencia preclínica",
                     chapter_tag="1 / 6")

    stats = [
        ("280+ M", "personas viven\ncon depresión (OMS)", ACCENT),
        ("8,000+", "artículos publicados\ncon FST desde 2000", BLUE),
        ("Desde 1977", "paradigma de\nPorsolt et al.", MIDNIGHT),
    ]
    card_w = 3.7; card_h = 1.55; gap = 0.4
    start_x = (SW - (3*card_w + 2*gap))/2
    for i,(big,small,c) in enumerate(stats):
        x = start_x + i*(card_w+gap)
        add_rect(s, x, 1.85, card_w, card_h, WHITE)
        add_rect(s, x, 1.85, 0.09, card_h, c)
        add_text(s, x+0.25, 1.98, card_w-0.3, 0.65, big,
                 size=32, bold=True, color=c, font="Georgia")
        add_text(s, x+0.25, 2.75, card_w-0.3, 0.7, small,
                 size=12, color=TEXT_DARK)

    add_rect(s, 0.8, 3.75, SW-1.6, 3.3, WHITE)
    add_rect(s, 0.8, 3.75, 0.09, 3.3, NAVY)
    add_text(s, 1.05, 3.95, 11, 0.5, "El nado forzado como modelo preclínico",
             size=18, bold=True, color=NAVY)
    bullets = [
        "Paradigma conductual estándar para evaluar potencial antidepresivo de nuevas moléculas.",
        "La rata flota con movimientos mínimos tras el estrés → desesperanza conductual.",
        "Validez predictiva: los fármacos aprobados clínicamente reducen la inmovilidad.",
        "Usado en neurociencia preclínica para distinguir mecanismos serotoninérgicos y noradrenérgicos.",
    ]
    add_bullets(s, 1.1, 4.55, 11.3, 2.3, bullets, size=14)
    footer_bar(s, 3, total)

# ==========================================================
# SLIDE 4 — Planteamiento del problema
# ==========================================================
def slide_problema(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Planteamiento del problema",
                     "Costos de tiempo, variabilidad y falta de registro centralizado",
                     chapter_tag="1 / 6")

    add_rect(s, 0.8, 1.7, 5.9, 5.3, WHITE)
    add_rect(s, 0.8, 1.7, 5.9, 0.7, MIDNIGHT)
    add_text(s, 0.95, 1.7, 5.7, 0.7, "Laboratorio de Neurociencia Conductual",
             size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    add_text(s, 0.95, 2.5, 5.6, 0.35, "ENMyH-IPN  ·  Dr. Sandino Reyes López",
             size=12, color=TEXT_MUTE, italic=True)
    add_text(s, 0.95, 2.95, 5.6, 0.4, "Protocolo actual",
             size=14, bold=True, color=NAVY)
    info = [
        "Videos con 4 ratas simultáneas (vista lateral).",
        "Sesión Día 1: 20 min (línea base).",
        "Sesión Día 2: 5 min (post-tratamiento).",
        "Cronómetro en mano, cuadro a cuadro.",
    ]
    add_bullets(s, 1.0, 3.4, 5.5, 2.3, info, size=12)

    add_text(s, 7.1, 1.7, 5.6, 0.45, "Cuatro problemas concretos",
             size=17, bold=True, color=NAVY, font="Georgia")
    problemas = [
        ("A", "Tiempo",          "Horas por experimento completo."),
        ("B", "Variabilidad",    "Criterios diferentes entre evaluadores."),
        ("C", "Sin centralizar", "Resultados en hojas sueltas y notas."),
        ("D", "Sin granularidad","Desglose por minuto inviable a mano."),
    ]
    card_w = 5.6; card_h = 1.1
    for i,(ic, t, d) in enumerate(problemas):
        y = 2.25 + i*1.2
        add_rect(s, 7.1, y, card_w, card_h, WHITE)
        add_rect(s, 7.1, y, 0.09, card_h, WARN)
        add_oval(s, 7.3, y+0.2, 0.7, 0.7, WARN)
        add_text(s, 7.3, y+0.2, 0.7, 0.7, ic, size=22, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, 8.1, y+0.12, 4.4, 0.42, t, size=15, bold=True, color=NAVY)
        add_text(s, 8.1, y+0.55, 4.4, 0.5, d, size=12, color=TEXT_DARK)
    footer_bar(s, 4, total)

# ==========================================================
# SLIDE 5 — Propuesta de solución
# ==========================================================
def slide_propuesta(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Propuesta de solución",
                     "Sistema web asíncrono con pipeline de IA para tres conductas",
                     chapter_tag="1 / 6")

    components = [
        ("Pipeline de análisis de video",
         "Detecta cada rata en su cilindro y clasifica cuadro a cuadro su conducta.",
         ["Modelo de fondo + Otsu + morfología",
          "Detección con YOLOv8 (rat.pt)",
          "Seguimiento con ByteTrack + Kalman",
          "Clasificación con ResNet-18 (transfer learning)",
          "Respaldo: CSRT / KCF si falta GPU"],
         NAVY, "01"),
        ("Plataforma web",
         "Un investigador sube videos desde el navegador y recibe reportes cuantitativos.",
         ["Carga de videos Día 1 / Día 2",
          "Dashboard con estado y progreso en tiempo real",
          "Resultados por animal, sesión y minuto",
          "Reportes descargables en PDF y CSV",
          "Gestión de usuarios y roles"],
         BLUE, "02"),
    ]
    card_w = 5.9; card_h = 5.1
    start_x = (SW - (2*card_w + 0.5))/2
    for i,(t, d, items, c, num) in enumerate(components):
        x = start_x + i*(card_w+0.5)
        add_rect(s, x, 1.85, card_w, card_h, WHITE)
        add_rect(s, x, 1.85, card_w, 1.15, c)
        add_text(s, x+0.25, 1.95, 1.5, 1.0, num,
                 size=44, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, x+1.5, 2.0, card_w-1.7, 0.55, t,
                 size=17, bold=True, color=WHITE, anchor=MSO_ANCHOR.TOP, font="Georgia")
        add_text(s, x+1.5, 2.55, card_w-1.7, 0.4, d,
                 size=11, color=ICE, italic=True)
        add_bullets(s, x+0.25, 3.25, card_w-0.45, 3.6, items, size=12, bullet_color=c)
    footer_bar(s, 5, total)

# ==========================================================
# SLIDE 6 — Objetivos
# ==========================================================
def slide_objetivos(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Objetivos",
                     "General y 6 específicos técnicos",
                     chapter_tag="2 / 6")

    add_rect(s, 0.8, 1.70, SW-1.6, 1.20, NAVY)
    add_rect(s, 0.8, 1.70, 0.12, 1.20, ACCENT)
    add_text(s, 1.05, 1.78, SW-1.9, 0.38, "Objetivo general",
             size=13, bold=True, color=ACCENT, italic=True)
    add_text(s, 1.05, 2.15, SW-1.9, 0.75,
             "Diseñar, desarrollar y validar un prototipo de sistema web asistido por IA que automatice la detección y cuantificación de las conductas del FST en ratas, con resultados comparables a la evaluación manual y cumpliendo los atributos de calidad de la ISO/IEC 25010.",
             size=13, color=WHITE)

    add_text(s, 0.8, 3.05, SW-1.6, 0.35, "Objetivos específicos",
             size=15, bold=True, color=NAVY)

    especs = [
        ("1", "Revisar trabajos y herramientas de análisis automatizado del FST."),
        ("2", "Documentar el protocolo experimental y las condiciones de grabación."),
        ("3", "Diseñar la arquitectura conforme a ISO/IEC/IEEE 12207."),
        ("4", "Implementar el pipeline de visión + clasificador supervisado."),
        ("5", "Validar con κ de Cohen y MAE contra el gold standard."),
        ("6", "Cuantificar la reducción de tiempo vs. análisis manual."),
    ]
    col_w = 4.0; row_h = 1.45; gap_x = 0.15; gap_y = 0.18
    start_x = (SW - (3*col_w + 2*gap_x))/2
    for i,(n,t) in enumerate(especs):
        col = i%3; row = i//3
        x = start_x + col*(col_w+gap_x)
        y = 3.50 + row*(row_h+gap_y)
        add_rect(s, x, y, col_w, row_h, WHITE)
        add_rect(s, x, y, 0.09, row_h, BLUE)
        add_oval(s, x+0.25, y+0.25, 0.55, 0.55, BLUE)
        add_text(s, x+0.25, y+0.25, 0.55, 0.55, n, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, x+0.95, y+0.25, col_w-1.1, row_h-0.5, t,
                 size=12, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
    footer_bar(s, 6, total)

# ==========================================================
# SLIDE 7 — Justificación y Alcance
# ==========================================================
def slide_justificacion_alcance(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Justificación y Alcance",
                     "Cumplimiento ISO/IEC 25010:2023 · delimitaciones del prototipo",
                     chapter_tag="2 / 6")

    # Justificación — banda superior
    add_rect(s, 0.8, 1.75, SW-1.6, 1.55, WHITE)
    add_rect(s, 0.8, 1.75, 0.12, 1.55, ACCENT)
    add_text(s, 1.05, 1.85, 11.5, 0.4, "Justificación", size=14, bold=True, color=ACCENT)
    justif = [
        "Reduce horas de análisis manual y homogeniza criterios entre evaluadores.",
        "Código abierto desplegable con Docker (ISO/IEC/IEEE 12207:2017).",
        "Clasificador modular: base para otros paradigmas conductuales.",
        "Atributos de calidad verificables bajo ISO/IEC 25010:2023 (κ, MAE, tiempo, usabilidad).",
    ]
    add_bullets(s, 1.05, 2.2, 11.5, 1.1, justif, size=11)

    # Dentro / Fuera
    incluye = [
        "Análisis de videos FST del laboratorio (vista lateral, 4 cilindros).",
        "Clasificación de 3 conductas (nado, inmovilidad, escalamiento).",
        "Fallback: si no se alcanza κ>0.80, nado+escalamiento = 'conducta activa'.",
        "Métricas por animal y sesión; desglose por minuto (opcional).",
        "Comparación Día 1 vs. Día 2; reportes PDF/CSV.",
        "Validación contra gold standard (κ de Cohen y MAE).",
    ]
    fuera = [
        "Análisis en tiempo real sobre video en vivo.",
        "Corrección automática de videos con mala iluminación.",
        "Intervención con animales (lo hace el laboratorio bajo NOM-062).",
        "Despliegue institucional definitivo (prototipo funcional).",
        "Otros paradigmas (laberinto en cruz, campo abierto, etc.).",
    ]
    add_rect(s, 0.8, 3.5, 6.0, 3.5, WHITE)
    add_rect(s, 0.8, 3.5, 6.0, 0.55, ACCENT)
    add_text(s, 0.8, 3.5, 6.0, 0.55, "   ✓  Dentro del alcance",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    add_bullets(s, 1.0, 4.15, 5.7, 2.7, incluye, size=11, bullet_color=ACCENT, marker="✓")

    add_rect(s, 7.0, 3.5, 5.5, 3.5, WHITE)
    add_rect(s, 7.0, 3.5, 5.5, 0.55, WARN)
    add_text(s, 7.0, 3.5, 5.5, 0.55, "   ✕  Fuera del alcance",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    add_bullets(s, 7.2, 4.15, 5.2, 2.7, fuera, size=11, bullet_color=WARN, marker="✕")

    footer_bar(s, 7, total)

# ==========================================================
# SLIDE 8 — Estado del Arte
# ==========================================================
def slide_estado_arte(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Estado del Arte",
                     "EthoVision XT, ezTrack y la brecha abierta",
                     chapter_tag="3 / 6")

    headers = ["Sistema", "Esp. FST", "3 cond.", "DL", "Web", "Abierto", "MOT"]
    rows = [
        ["EthoVision XT",     "Parc.", "No",  "No",  "No",  "No",  "No"],
        ["ANY-maze",          "Parc.", "No",  "No",  "No",  "No",  "No"],
        ["ezTrack",           "Parc.", "No",  "No",  "No",  "Sí",  "No"],
        ["DBscorer",          "Sí",    "No",  "No",  "No",  "Sí",  "No"],
        ["DSAAN (2024)",      "Sí",    "No",  "Sí",  "No",  "No",  "No"],
        ["3D-RCNN (2025)",    "Sí",    "Sí",  "Sí",  "No",  "No",  "No"],
        ["DeepLabCut + clf",  "Parc.", "Parc.","Sí", "No",  "Sí",  "No"],
        ["Propuesta",         "Sí",    "Sí",  "Sí",  "Sí",  "Sí",  "ByteTrack"],
    ]
    x0 = 0.8; y0 = 1.85
    col_ws = [2.6, 1.3, 1.3, 1.0, 1.0, 1.3, 1.95]
    row_h = 0.44

    x = x0
    for i,h in enumerate(headers):
        add_rect(s, x, y0, col_ws[i], row_h, NAVY)
        add_text(s, x, y0, col_ws[i], row_h, h, size=12, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_ws[i]

    for r_idx, row in enumerate(rows):
        y = y0 + (r_idx+1)*row_h
        is_prop = (r_idx == len(rows)-1)
        fill = ACCENT if is_prop else (WHITE if r_idx%2==0 else ICE)
        text_color = WHITE if is_prop else TEXT_DARK
        x = x0
        for i,cell in enumerate(row):
            add_rect(s, x, y, col_ws[i], row_h, fill)
            add_text(s, x, y, col_ws[i], row_h, cell,
                     size=11, bold=(i==0 or is_prop),
                     color=text_color,
                     align=PP_ALIGN.CENTER if i>0 else PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += col_ws[i]

    y_concl = y0 + (len(rows)+1)*row_h + 0.25
    add_rect(s, 0.8, y_concl, SW-1.6, 0.9, MIDNIGHT)
    add_rect(s, 0.8, y_concl, 0.09, 0.9, ACCENT)
    add_text(s, 1.05, y_concl+0.12, SW-1.9, 0.35, "La brecha",
             size=13, bold=True, color=ACCENT)
    add_text(s, 1.05, y_concl+0.42, SW-1.9, 0.4,
             "Las comerciales no son abiertas ni distinguen las 3 conductas. Las abiertas no usan DL ni tienen web. Los trabajos DL recientes no se distribuyen como sistema desplegable.",
             size=11, color=WHITE)
    footer_bar(s, 8, total)

# ==========================================================
# SLIDE 9 — Marco Teórico Conductual
# ==========================================================
def slide_marco_conductual(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Marco Teórico Conductual",
                     "Paradigma FST: tres conductas mutuamente excluyentes",
                     chapter_tag="3 / 6")

    add_text(s, 0.8, 1.7, SW-1.6, 0.5,
             "Cada conducta refleja un sistema neurotransmisor distinto. Distinguirlas permite inferir sobre qué mecanismo actúa el fármaco evaluado.",
             size=13, color=TEXT_DARK, italic=True)

    conductas = [
        ("Inmovilidad",
         "El animal flota con lo mínimo para no hundirse.",
         "Indicador principal de desesperanza conductual.",
         NAVY),
        ("Nado activo",
         "Paladas que desplazan al animal horizontalmente.",
         "Asociado al sistema serotoninérgico (ISRS, fluoxetina).",
         BLUE),
        ("Intentos de escape",
         "Movimientos rápidos hacia las paredes del cilindro.",
         "Asociado al sistema noradrenérgico (desipramina).",
         ACCENT),
    ]
    card_w = 3.9; card_h = 3.7; gap = 0.3
    start_x = (SW - (3*card_w + 2*gap))/2
    for i,(t,d,extra,c) in enumerate(conductas):
        x = start_x + i*(card_w+gap)
        add_rect(s, x, 2.45, card_w, card_h, WHITE)
        add_rect(s, x, 2.45, card_w, 0.6, c)
        add_text(s, x, 2.45, card_w, 0.6, t, size=19, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_oval(s, x+card_w/2-0.55, 3.25, 1.1, 1.1, ICE)
        add_text(s, x+card_w/2-0.55, 3.25, 1.1, 1.1, ["I","II","III"][i],
                 size=36, bold=True, color=c, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, x+0.25, 4.55, card_w-0.5, 0.7, d,
                 size=12, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x+0.25, 5.25, card_w-0.5, 1.6, extra,
                 size=11, color=TEXT_MUTE, align=PP_ALIGN.CENTER, italic=True)
    footer_bar(s, 9, total)

# ==========================================================
# SLIDE 10 — Marco Teórico Visión + Aprendizaje Profundo
# ==========================================================
def slide_marco_vision_dl(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Marco Teórico de Visión y Aprendizaje Profundo",
                     "Sustracción de fondo · ByteTrack · ResNet-18 · Transfer learning",
                     chapter_tag="3 / 6")

    # 4 bloques en 2x2
    bloques = [
        ("Sustracción de fondo",
         "Modelo del fondo por mediana pixel a pixel de los primeros N cuadros. CLAHE corrige contraste; morfología (apertura) elimina ruido tras el umbral.",
         NAVY),
        ("ByteTrack",
         "Rastreo multi-objeto con filtro de Kalman y asociación en dos pasos: detecciones de alta confianza (s ≥ τ_high) y luego las de baja (τ_low ≤ s < τ_high) para recuperar trayectorias en oclusiones.",
         BLUE),
        ("ResNet-18",
         "Red residual con conexiones de salto que aprenden F(x)=H(x)−x. El gradiente fluye sin atenuación, habilita redes profundas. ResNet-18 balancea precisión e inferencia rápida.",
         MIDNIGHT),
        ("Transfer learning",
         "Partimos de ImageNet (≈1M imágenes, 1000 clases); congelamos las capas iniciales y reentrenamos solo la capa final con 3 salidas (una por conducta) usando cuadros etiquetados en BORIS.",
         ACCENT),
    ]
    card_w = 5.9; card_h = 2.35; gap_x = 0.3; gap_y = 0.2
    start_x = (SW - (2*card_w + gap_x))/2
    for i,(t, d, c) in enumerate(bloques):
        col = i%2; row = i//2
        x = start_x + col*(card_w+gap_x)
        y = 1.75 + row*(card_h+gap_y)
        add_rect(s, x, y, card_w, card_h, WHITE)
        add_rect(s, x, y, 0.12, card_h, c)
        add_rect(s, x+0.12, y, card_w-0.12, 0.5, c)
        add_text(s, x+0.3, y, card_w-0.35, 0.5, t,
                 size=14, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, x+0.3, y+0.65, card_w-0.45, card_h-0.8, d,
                 size=11, color=TEXT_DARK)
    footer_bar(s, 10, total)

# ==========================================================
# SLIDE 11 — Metodología (Scrum)
# ==========================================================
def slide_metodologia(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Metodología",
                     "Enfoque ágil Scrum · distribución por integrante",
                     chapter_tag="4 / 6")

    # Columna izquierda: fases Scrum
    add_rect(s, 0.8, 1.75, 6.0, 5.3, WHITE)
    add_rect(s, 0.8, 1.75, 6.0, 0.6, NAVY)
    add_text(s, 0.8, 1.75, 6.0, 0.6, "  Flujo iterativo por sprints",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    fases = [
        ("Backlog", "RF/RNF priorizados por riesgo técnico."),
        ("Sprint planning", "Incremento entregable cada 2 semanas."),
        ("Daily", "Sincronización breve; desbloqueo temprano."),
        ("Review", "Demostración con el director y laboratorio."),
        ("Retrospectiva", "Ajuste de criterios y estimación."),
    ]
    y = 2.5
    for t, d in fases:
        add_rect(s, 1.0, y, 5.6, 0.8, BG_LIGHT)
        add_rect(s, 1.0, y, 0.1, 0.8, ACCENT)
        add_text(s, 1.2, y+0.08, 5.3, 0.32, t, size=13, bold=True, color=NAVY)
        add_text(s, 1.2, y+0.4, 5.3, 0.38, d, size=10, color=TEXT_DARK)
        y += 0.9

    # Columna derecha: distribución de actividades
    add_rect(s, 7.0, 1.75, 5.6, 5.3, WHITE)
    add_rect(s, 7.0, 1.75, 5.6, 0.6, ACCENT)
    add_text(s, 7.0, 1.75, 5.6, 0.6, "  Distribución por integrante",
             size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    roles = [
        ("Frausto Robles Ángel Ali",
         "Pipeline de visión, detección YOLOv8, ByteTrack, clasificador ResNet-18, pruebas de κ/MAE."),
        ("Rodríguez Verdín Sandoval Vanesa",
         "Plataforma web Flask + React, modelado de BD, worker asíncrono, dashboard y reportes."),
        ("Trabajo conjunto",
         "Análisis de requisitos, diseño de arquitectura, integración end-to-end y documentación."),
    ]
    y = 2.5
    for t, d in roles:
        add_rect(s, 7.2, y, 5.2, 1.45, BG_LIGHT)
        add_rect(s, 7.2, y, 0.1, 1.45, BLUE)
        add_text(s, 7.4, y+0.08, 4.9, 0.4, t, size=13, bold=True, color=NAVY)
        add_text(s, 7.4, y+0.48, 4.9, 0.95, d, size=11, color=TEXT_DARK)
        y += 1.55
    footer_bar(s, 11, total)

# ==========================================================
# SLIDE 12 — Requerimientos Funcionales (RF-01 a RF-29)
# ==========================================================
def slide_rf(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Requerimientos funcionales",
                     "29 RF organizados por módulos del sistema",
                     chapter_tag="4 / 6")

    modulos = [
        ("Autenticación", "RF-01 → RF-04", NAVY, [
            "Login con @ipn.mx · logout JWT",
            "Recuperación y primer cambio de contraseña",
            "Roles: Investigador / Administrador",
        ]),
        ("Gestión de usuarios", "RF-05 → RF-07", MIDNIGHT, [
            "Admin: crear, modificar, desactivar cuentas",
            "Investigador: actualizar perfil propio",
        ]),
        ("Carga de video", "RF-08 → RF-12", BLUE, [
            "Solo .mp4 · hasta 2 videos (DAY1/DAY2)",
            "Validación cliente+servidor · metadatos",
            "Encolado automático tras carga exitosa",
        ]),
        ("Análisis conductual", "RF-13 → RF-17", ACCENT, [
            "Pipeline 4 etapas · progreso en tiempo real",
            "Detección de cilindros (conf. ≥ 0.70)",
            "Clasificación frame-a-frame mutuamente excluyente",
            "Reporte de error con PDF de diagnóstico",
        ]),
        ("Resultados", "RF-18 → RF-21", BLUE, [
            "Tiempo y % por conducta y animal",
            "Desglose por minuto · comparación DAY1 vs DAY2",
            "Exportación PDF y CSV",
        ]),
        ("Dashboard y administración", "RF-22 → RF-29", MIDNIGHT, [
            "Listado de experimentos y estado",
            "Avisos de expiración · borrado a 30 días",
            "Admin: monitorear disco, cola, análisis activos",
        ]),
    ]
    card_w = 4.0; card_h = 2.5; gap_x = 0.18; gap_y = 0.18
    start_x = (SW - (3*card_w + 2*gap_x))/2
    for i,(t, rng, c, items) in enumerate(modulos):
        col = i%3; row = i//3
        x = start_x + col*(card_w+gap_x)
        y = 1.8 + row*(card_h+gap_y)
        add_rect(s, x, y, card_w, card_h, WHITE)
        add_rect(s, x, y, card_w, 0.5, c)
        add_text(s, x+0.15, y, card_w*0.65, 0.5, t,
                 size=12, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, x+card_w*0.55, y, card_w*0.43, 0.5, rng,
                 size=10, bold=True, color=ICE,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_bullets(s, x+0.15, y+0.55, card_w-0.25, card_h-0.6, items,
                    size=9, bullet_color=c, line_sp=1.05)
    footer_bar(s, 12, total)

# ==========================================================
# SLIDE 13 — RNF + Reglas de negocio
# ==========================================================
def slide_rnf_reglas(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Requerimientos no funcionales y reglas de negocio",
                     "Atributos ISO 25010 · lógica del protocolo experimental",
                     chapter_tag="4 / 6")

    # Columna izquierda: RNF
    add_rect(s, 0.8, 1.75, 6.0, 5.3, WHITE)
    add_rect(s, 0.8, 1.75, 6.0, 0.55, NAVY)
    add_text(s, 0.8, 1.75, 6.0, 0.55, "  RNF · 13 atributos de calidad",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    rnf = [
        ("Rendimiento", "Navegación < 3 s (p95) · umbral de análisis por minuto de video"),
        ("Disponibilidad", "95 % en horario laboral L-V 8–20h"),
        ("Seguridad", "JWT · bcrypt+salt · aislamiento por usuario"),
        ("Usabilidad", "Investigador sin experiencia completa flujo"),
        ("Compatibilidad", "Chrome, Firefox y Edge actuales"),
        ("Portabilidad", "Despliegue único con Docker Compose"),
        ("Mantenibilidad", "PEP-8 · ESLint · documentación por módulo"),
    ]
    y = 2.4
    for t, d in rnf:
        add_text(s, 1.0, y, 5.8, 0.3, t, size=11, bold=True, color=ACCENT)
        add_text(s, 1.0, y+0.28, 5.8, 0.35, d, size=10, color=TEXT_DARK)
        y += 0.62

    # Columna derecha: Reglas de negocio
    add_rect(s, 7.0, 1.75, 5.6, 5.3, WHITE)
    add_rect(s, 7.0, 1.75, 5.6, 0.55, ACCENT)
    add_text(s, 7.0, 1.75, 5.6, 0.55, "  Reglas de negocio · 13 RN",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    rn = [
        "RN-01 · Error de credenciales genérico.",
        "RN-02 · Video único por experimento.",
        "RN-03 · Máximo 2 videos (DAY1 + DAY2).",
        "RN-04 · Análisis iniciado por el sistema.",
        "RN-05 · Video borrado a 30 días + aviso 7 días.",
        "RN-06 · Resultados persistentes indefinidamente.",
        "RN-07 · Una conducta por cuadro por animal.",
        "RN-08 · Investigador ve lo propio · Admin ve todo.",
        "RN-09 · Eliminación permanente con confirmación.",
        "RN-10 · Rechazo de no-MP4 en cliente.",
        "RN-11 · Detección de cilindros con conf. ≥ 0.70.",
        "RN-12 · DAY1/DAY2 comparables solo si ambos exitosos.",
        "RN-13 · Si κ ≤ 0.80 → fallback 'conducta activa'.",
    ]
    add_bullets(s, 7.2, 2.4, 5.3, 4.5, rn, size=9.5, bullet_color=WARN, line_sp=1.1)
    footer_bar(s, 13, total)

# ==========================================================
# SLIDE 14 — Análisis de Riesgos y Factibilidad
# ==========================================================
def slide_riesgos_factibilidad(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Análisis de Riesgos y Factibilidad",
                     "Mitigación de calidad · factibilidad técnica y económica",
                     chapter_tag="5 / 6")

    # Riesgos — tabla compacta arriba
    add_text(s, 0.8, 1.65, 11, 0.35, "Riesgos principales",
             size=14, bold=True, color=NAVY)
    riesgos = [
        ("R-01", "Videos con reflejos o pelaje claro (conf. < 0.70).",
         "Alto", "CLAHE · reporte de diagnóstico (RN-11)."),
        ("R-02", "Dataset insuficiente para κ > 0.80.",
         "Alto", "BORIS · data augmentation · fallback."),
        ("R-03", "Laboratorio no entrega registros a tiempo.",
         "Medio", "Calendario pactado · anotación anticipada."),
        ("R-04", "Tiempo de análisis excede lo aceptable.",
         "Medio", "Medir por etapa · GPU opcional."),
        ("R-05", "Integración cola asíncrona compleja.",
         "Medio", "Contratos API · pruebas end-to-end."),
    ]
    y = 2.00; row_h = 0.50
    cols = [("ID",1.0), ("Descripción",6.8), ("Nivel",1.2), ("Mitigación",2.6)]
    x = 0.8
    add_rect(s, 0.8, y, 11.6, row_h, NAVY)
    for t,w in cols:
        add_text(s, x, y, w, row_h, t, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += w
    y += row_h
    for i,(rid, desc, lvl, mit) in enumerate(riesgos):
        fill = WHITE if i%2==0 else ICE
        add_rect(s, 0.8, y, 11.6, row_h, fill)
        add_text(s, 0.8, y, 1.0, row_h, rid, size=11, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.8, y, 6.7, row_h, desc, size=10, color=TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        lvl_color = WARN if lvl=="Alto" else BLUE
        add_rounded(s, 8.7, y+0.08, 0.95, 0.34, lvl_color)
        add_text(s, 8.7, y+0.08, 0.95, 0.34, lvl, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 9.8, y, 2.55, row_h, mit, size=9, color=TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += row_h

    # Factibilidad — 3 tarjetas inferiores
    add_text(s, 0.8, 5.15, 11, 0.32, "Factibilidad",
             size=14, bold=True, color=NAVY)
    fact = [
        ("Técnica", "Todo el stack es código abierto y maduro: OpenCV, YOLOv8, PyTorch, Flask, PostgreSQL.", NAVY),
        ("Económica", "Costo de licencias: $0. Solo requiere máquina con CPU moderna; GPU opcional.", ACCENT),
        ("Operativa", "Despliegue con Docker Compose en el equipo del laboratorio; no requiere infraestructura nueva.", BLUE),
    ]
    card_w = 4.0; card_h = 1.30; gap = 0.2
    start_x = (SW - (3*card_w + 2*gap))/2
    for i,(t,d,c) in enumerate(fact):
        x = start_x + i*(card_w+gap)
        add_rect(s, x, 5.50, card_w, card_h, WHITE)
        add_rect(s, x, 5.50, 0.1, card_h, c)
        add_text(s, x+0.2, 5.58, card_w-0.3, 0.32, t, size=13, bold=True, color=c)
        add_text(s, x+0.2, 5.90, card_w-0.3, 0.85, d, size=10, color=TEXT_DARK)
    footer_bar(s, 14, total)

# ==========================================================
# SLIDE 15 — Diseño de Arquitectura (cliente-servidor + REST + workers)
# ==========================================================
def slide_arquitectura(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Diseño de Arquitectura",
                     "Cliente-servidor · API REST · procesamiento asíncrono",
                     chapter_tag="5 / 6")

    # Diagrama a la izquierda (image1.jpg con paquetes)
    img = os.path.join(DIAG, "image1.jpg")
    s.shapes.add_picture(img, Inches(0.8), Inches(1.7), height=Inches(4.3))

    # Panel derecho: capas
    add_rect(s, 8.6, 1.7, 4.0, 4.3, WHITE)
    add_rect(s, 8.6, 1.7, 0.1, 4.3, NAVY)
    add_text(s, 8.8, 1.8, 3.7, 0.35, "Cinco paquetes", size=13, bold=True, color=NAVY)
    add_bullets(s, 8.8, 2.15, 3.7, 2.3,
        ["P1 · Autenticación",
         "P2 · Carga de video",
         "P3 · Pipeline de IA",
         "P4 · Resultados y reportes",
         "P5 · Dashboard y admin"], size=10)
    add_text(s, 8.8, 4.55, 3.7, 0.35, "Actores", size=13, bold=True, color=NAVY)
    add_bullets(s, 8.8, 4.9, 3.7, 1.1,
        ["Investigador (principal)",
         "Administrador",
         "Worker (sistema)"], size=10, bullet_color=BLUE)

    # Tira inferior: flujo cliente-servidor-worker
    add_rect(s, 0.8, 6.2, SW-1.6, 0.9, WHITE)
    add_rect(s, 0.8, 6.2, 0.1, 0.9, ACCENT)
    capas = [
        ("React + Vite", "frontend en el navegador", BLUE),
        ("Flask + JWT",  "API REST · valida · encola", NAVY),
        ("PostgreSQL",   "persistencia ACID", MIDNIGHT),
        ("Worker",       "pipeline IA asíncrono", ACCENT),
    ]
    box_w = 2.75; box_h = 0.68; gx = 0.08
    start_x = 1.1
    for i,(t, d, c) in enumerate(capas):
        x = start_x + i*(box_w+gx)
        add_rect(s, x, 6.31, box_w, box_h, BG_LIGHT)
        add_rect(s, x, 6.31, box_w, 0.26, c)
        add_text(s, x, 6.31, box_w, 0.26, t, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, 6.57, box_w, 0.42, d, size=9, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, italic=True)
        if i < len(capas)-1:
            add_text(s, x+box_w-0.04, 6.35, gx+0.08, 0.6, "→",
                     size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
    footer_bar(s, 15, total)

# ==========================================================
# SLIDE 16 — Diseño del Pipeline de IA
# ==========================================================
def slide_pipeline_ia(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Diseño del Pipeline de IA",
                     "Preprocesamiento · detección · seguimiento · clasificación",
                     chapter_tag="5 / 6")

    # Diagrama del paquete 3 (image4.png)
    img = os.path.join(DIAG, "image4.png")
    s.shapes.add_picture(img, Inches(0.85), Inches(1.75), height=Inches(4.9))

    # 4 etapas a la derecha
    etapas = [
        ("1", "Preprocesamiento",
         "Decodificación + CLAHE adaptativo si el contraste cae bajo umbral; morfología de apertura para ruido.",
         NAVY),
        ("2", "Detección de cilindros",
         "Localiza los 4 ROIs con confianza ≥ 0.70 (RN-11); si falla, se genera reporte de diagnóstico.",
         BLUE),
        ("3", "Seguimiento",
         "YOLOv8 + ByteTrack con filtro de Kalman; respaldo CSRT/KCF si no hay GPU.",
         MIDNIGHT),
        ("4", "Clasificación",
         "ResNet-18 con transfer learning (ImageNet → 3 conductas); ResNet-50 como alternativa.",
         ACCENT),
    ]
    top = 1.85; step_h = 1.2
    for i,(n,t,d,c) in enumerate(etapas):
        y = top + i*step_h
        add_rect(s, 7.3, y, 5.5, step_h-0.1, WHITE)
        add_rect(s, 7.3, y, 0.1, step_h-0.1, c)
        add_oval(s, 7.5, y+0.2, 0.7, 0.7, c)
        add_text(s, 7.5, y+0.2, 0.7, 0.7, n, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, 8.3, y+0.12, 4.4, 0.4, t, size=13, bold=True, color=NAVY)
        add_text(s, 8.3, y+0.5, 4.4, 0.6, d, size=10, color=TEXT_DARK)
    footer_bar(s, 16, total)

# ==========================================================
# SLIDE 17 — Diseño de Base de Datos (ER)
# ==========================================================
def slide_bd(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Diseño de la Base de Datos",
                     "Modelo entidad–relación · persistencia en PostgreSQL",
                     chapter_tag="5 / 6")

    # Diagrama ER al centro (imagen portrait 1200x1846) — maximizar altura
    img = os.path.join(DIAG, "er.png")
    pic_h = 5.40
    pic_w = pic_h * (1200/1846)  # ≈ 3.51
    pic_x = (SW - pic_w) / 2
    s.shapes.add_picture(img, Inches(pic_x), Inches(1.68),
                         width=Inches(pic_w), height=Inches(pic_h))

    # Panel izquierdo — entidades
    add_rect(s, 0.55, 1.68, 3.05, 5.40, WHITE)
    add_rect(s, 0.55, 1.68, 0.09, 5.40, NAVY)
    add_text(s, 0.75, 1.80, 2.80, 0.36, "Entidades principales",
             size=13, bold=True, color=NAVY)
    ents = [
        ("USERS",       "Rol (investigador / admin)"),
        ("EXPERIMENTS", "Nombre, fecha, tratamiento"),
        ("VIDEOS",      "DAY1 / DAY2, deletion_date"),
        ("JOBS",        "Estado, etapa, progress_pct"),
        ("ANIMALS",     "ROI por sujeto (1–4)"),
        ("BEHAVIOR_RESULTS", "swim · immobile · escape"),
        ("REPORTS",     "PDF / CSV"),
        ("NOTIFICATIONS", "análisis, expiración, error"),
    ]
    y = 2.22
    for name, desc in ents:
        add_text(s, 0.80, y, 2.75, 0.28, name,
                 size=10, bold=True, color=ACCENT)
        add_text(s, 0.80, y+0.24, 2.75, 0.26, desc, size=9, color=TEXT_DARK)
        y += 0.56

    # Panel derecho — claves
    add_rect(s, SW-3.00, 1.68, 2.45, 5.40, WHITE)
    add_rect(s, SW-3.00, 1.68, 0.09, 5.40, BLUE)
    add_text(s, SW-2.80, 1.80, 2.25, 0.36, "Claves del diseño",
             size=13, bold=True, color=NAVY)
    notas = [
        "PostgreSQL con ACID.",
        "SQLAlchemy como ORM.",
        "Borrar experimento elimina todo lo asociado (composición).",
        "Resultados sobreviven al borrado del video (RN-06).",
        "Jobs ligan video↔pipeline y conservan progreso por etapa.",
        "Notifications: análisis listo, expiración 30 d, error.",
    ]
    add_bullets(s, SW-2.80, 2.22, 2.30, 4.60, notas, size=10)
    footer_bar(s, 17, total)

# ==========================================================
# SLIDE 18 — Vistas del Prototipo
# ==========================================================
def slide_vistas(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Vistas del Prototipo",
                     "Dashboard · creación de experimento · resultados",
                     chapter_tag="5 / 6")

    vistas = [
        ("Login",
         "Acceso con correo @ipn.mx",
         ["Campo email + contraseña",
          "Mensaje de error genérico (RN-01)",
          "Enlace de recuperación"],
         NAVY),
        ("Dashboard del investigador",
         "Lista de experimentos y estado",
         ["Progreso por análisis activo",
          "Aviso de video <7 días a expirar",
          "Acceso directo a cada experimento"],
         BLUE),
        ("Creación de experimento",
         "Formulario en tres pasos",
         ["Metadatos (nombre, fecha, tratamiento)",
          "Asignación de ROIs a animales",
          "Carga de videos DAY1 / DAY2"],
         MIDNIGHT),
        ("Progreso del análisis",
         "Barra en tiempo real",
         ["4 etapas visibles del pipeline",
          "Porcentaje por etapa",
          "Manejo de error con reporte"],
         WARN),
        ("Resultados del experimento",
         "Cuatro pestañas",
         ["Resumen por conducta",
          "Detalle por animal",
          "Desglose por minuto · DAY1 vs DAY2"],
         ACCENT),
        ("Panel de administración",
         "Métricas y gestión",
         ["Uso de disco · cola · análisis activos",
          "Gestión de usuarios y roles",
          "Alertas al 80 % de capacidad"],
         MIDNIGHT),
    ]
    card_w = 4.0; card_h = 2.5; gap_x = 0.18; gap_y = 0.18
    start_x = (SW - (3*card_w + 2*gap_x))/2
    for i,(t, sub, items, c) in enumerate(vistas):
        col = i%3; row = i//3
        x = start_x + col*(card_w+gap_x)
        y = 1.8 + row*(card_h+gap_y)
        add_rect(s, x, y, card_w, card_h, WHITE)
        add_rect(s, x, y, card_w, 0.5, c)
        add_text(s, x+0.2, y, card_w-0.3, 0.5, t,
                 size=12, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
        add_text(s, x+0.15, y+0.55, card_w-0.25, 0.3, sub,
                 size=10, color=TEXT_MUTE, italic=True)
        add_bullets(s, x+0.15, y+0.9, card_w-0.25, card_h-1.0, items,
                    size=9, bullet_color=c, line_sp=1.1)
    footer_bar(s, 18, total)

# ==========================================================
# SLIDE 19 — Validación y Métricas
# ==========================================================
def slide_metricas(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Validación y Métricas",
                     "κ de Cohen · F1-score · Error Absoluto Medio",
                     chapter_tag="6 / 6")

    atributos = [
        ("κ > 0.80", "κ de Cohen",
         "Concordancia inter-evaluador vs anotaciones expertas (gold standard con BORIS).",
         NAVY),
        ("F1", "F1-score por clase",
         "Balance precisión/recall por conducta; detecta sesgos de clase minoritaria.",
         BLUE),
        ("MAE", "Error absoluto medio",
         "Segundos por animal y conducta respecto al registro manual.",
         ACCENT),
    ]
    card_w = 4.0; card_h = 2.6; gap = 0.25
    start_x = (SW - (3*card_w + 2*gap))/2
    for i,(big, label, desc, c) in enumerate(atributos):
        x = start_x + i*(card_w+gap)
        add_rect(s, x, 1.85, card_w, card_h, WHITE)
        add_rect(s, x, 1.85, 0.09, card_h, c)
        add_text(s, x+0.25, 2.0, card_w-0.4, 0.9, big,
                 size=38, bold=True, color=c, font="Georgia")
        add_text(s, x+0.25, 2.95, card_w-0.4, 0.5, label,
                 size=13, bold=True, color=NAVY)
        add_text(s, x+0.25, 3.5, card_w-0.4, 0.9, desc,
                 size=11, color=TEXT_DARK)

    add_rect(s, 0.8, 4.7, SW-1.6, 2.4, WHITE)
    add_rect(s, 0.8, 4.7, 0.12, 2.4, MIDNIGHT)
    add_text(s, 1.1, 4.85, 11, 0.45, "Usability y Time behaviour",
             size=16, bold=True, color=NAVY, font="Georgia")
    add_text(s, 1.1, 5.3, 11, 0.4,
             "Un investigador sin experiencia en programación debe crear un experimento, subir el video, esperar el análisis y descargar el reporte sin ayuda técnica.",
             size=12, color=TEXT_DARK)
    add_text(s, 1.1, 6.0, 11, 0.4, "Criterios de aceptación", size=12, bold=True, color=NAVY)
    add_bullets(s, 1.1, 6.35, 11, 0.6,
        ["Tasa de completación ≥ 90 % sin asistencia · al menos 1 usuario real del laboratorio.",
         "Tiempo de procesamiento ≤ tiempo de video en equipo del laboratorio (sin GPU)."],
        size=11)
    footer_bar(s, 19, total)

# ==========================================================
# SLIDE 20 — Conclusiones y Trabajo Futuro
# ==========================================================
def slide_conclusiones(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Conclusiones y Trabajo Futuro",
                     "Entregables de TT1 · siguientes pasos en TT2",
                     chapter_tag="6 / 6")

    # Conclusiones
    add_rect(s, 0.8, 1.75, 6.0, 5.3, WHITE)
    add_rect(s, 0.8, 1.75, 6.0, 0.55, NAVY)
    add_text(s, 0.8, 1.75, 6.0, 0.55, "  Conclusiones de TT1",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    concl = [
        "Se documentó el protocolo FST del laboratorio y se cuantificó el costo del análisis manual.",
        "El estado del arte confirmó la brecha: no existe un sistema abierto, web y con 3 conductas.",
        "Se definieron 29 RF, 13 RNF y 13 reglas de negocio alineados a ISO/IEC 25010:2023.",
        "La arquitectura propuesta (React + Flask + PostgreSQL + worker) es factible con costo $0 en licencias.",
        "El pipeline de IA (YOLOv8 + ByteTrack + ResNet-18) es viable sin GPU con respaldo clásico.",
    ]
    add_bullets(s, 1.0, 2.45, 5.7, 4.5, concl, size=11, bullet_color=ACCENT, line_sp=1.2)

    # Trabajo futuro
    add_rect(s, 7.0, 1.75, 5.6, 5.3, WHITE)
    add_rect(s, 7.0, 1.75, 5.6, 0.55, ACCENT)
    add_text(s, 7.0, 1.75, 5.6, 0.55, "  Trabajo futuro · TT2",
             size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, font="Georgia")
    futuro = [
        ("Implementación", "Backend Flask con JWT, BD y API REST."),
        ("Frontend",       "React + Vite con dashboard y reportes."),
        ("Entrenamiento",  "Ampliar dataset anotado en BORIS y entrenar ResNet-18."),
        ("Validación",     "Medir κ, F1 y MAE contra el gold standard."),
        ("Entrega",        "Docker Compose, manuales y material de defensa."),
    ]
    y = 2.45
    for t, d in futuro:
        add_rect(s, 7.2, y, 5.2, 0.85, BG_LIGHT)
        add_rect(s, 7.2, y, 0.1, 0.85, BLUE)
        add_text(s, 7.4, y+0.08, 4.9, 0.3, t, size=12, bold=True, color=NAVY)
        add_text(s, 7.4, y+0.42, 4.9, 0.4, d, size=10, color=TEXT_DARK)
        y += 0.93
    footer_bar(s, 20, total)

# ==========================================================
# SLIDE 21 — Bibliografía
# ==========================================================
def slide_bibliografia(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, BG_LIGHT)
    slide_header_bar(s, "Bibliografía",
                     "Fuentes citadas conforme al formato institucional",
                     chapter_tag="6 / 6")

    refs = [
        "Porsolt, R. D., Le Pichon, M., & Jalfre, M. (1977). Depression: a new animal model sensitive to antidepressant treatments. Nature, 266(5604), 730-732.",
        "Zhang, Y. et al. (2022). ByteTrack: Multi-object tracking by associating every detection box. ECCV 2022, pp. 1-21.",
        "He, K., Zhang, X., Ren, S., & Sun, J. (2016). Deep residual learning for image recognition. CVPR, pp. 770-778.",
        "Jocher, G. et al. (2023). YOLOv8 by Ultralytics. Documentación oficial. Ultralytics.",
        "Zuiderveld, K. (1994). Contrast limited adaptive histogram equalization. Graphics Gems IV, pp. 474-485.",
        "ISO/IEC 25010:2023. Systems and software engineering — Systems and software Quality Requirements and Evaluation (SQuaRE).",
        "ISO/IEC/IEEE 12207:2017. Systems and software engineering — Software life cycle processes.",
        "Friard, O., & Gamba, M. (2016). BORIS: Behavioral Observation Research Interactive Software. Methods in Ecology and Evolution, 7(11), 1325-1330.",
        "Pennington, Z. T. et al. (2019). ezTrack: An open-source video analysis pipeline. Scientific Reports, 9(1), 19979.",
        "Organización Mundial de la Salud (2023). Depression and Other Common Mental Disorders: Global Health Estimates.",
        "Schwaber, K., & Sutherland, J. (2020). The Scrum Guide.",
        "NOM-062-ZOO-1999. Especificaciones técnicas para la producción, cuidado y uso de los animales de laboratorio. DOF, México.",
    ]
    # Dos columnas
    col_w = 5.9; gap = 0.3
    start_x = 0.8
    half = (len(refs)+1)//2
    left = refs[:half]; right = refs[half:]

    add_bullets(s, start_x, 1.8, col_w, 5.3, left, size=9.5,
                bullet_color=NAVY, line_sp=1.15)
    add_bullets(s, start_x+col_w+gap, 1.8, col_w, 5.3, right, size=9.5,
                bullet_color=NAVY, line_sp=1.15)
    footer_bar(s, 21, total)

# ==========================================================
# SLIDE 22 — Cierre
# ==========================================================
def slide_cierre(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, 0, SW, 0.18, ACCENT)
    add_rect(s, 0, SH-0.18, SW, 0.18, ACCENT)

    add_text(s, 0, 1.7, SW, 0.5,
             "TRABAJO TERMINAL  ·  2026-B066",
             size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s, 0.8, 2.25, SW-1.6, 1.4,
             "De horas frente al cronómetro\na minutos de análisis automatizado.",
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Georgia")

    add_text(s, 0.8, 3.9, SW-1.6, 0.6,
             "Un prototipo abierto, desplegable con Docker y validado contra el gold standard del laboratorio.",
             size=15, color=ICE, italic=True, align=PP_ALIGN.CENTER)

    add_line(s, 4.5, 4.85, 8.85, 4.85, ACCENT, weight=2)

    add_text(s, 0, 5.1, SW, 0.7, "Gracias",
             size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Georgia")
    add_text(s, 0, 6.0, SW, 0.45, "¿Preguntas?",
             size=18, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

    add_text(s, 0, 6.6, SW, 0.4,
             "Frausto Robles Ángel Ali   ·   Rodríguez Verdín Sandoval Vanesa",
             size=12, color=ICE, align=PP_ALIGN.CENTER)
    add_text(s, 0, 6.95, SW, 0.35,
             "ESCOM-IPN  ·  Mayo 2026",
             size=10, color=ICE, align=PP_ALIGN.CENTER, italic=True)

# ==========================================================
# ORQUESTA — 22 slides
# ==========================================================
TOTAL = 22
slide_portada()
slide_agenda(TOTAL)
slide_introduccion(TOTAL)
slide_problema(TOTAL)
slide_propuesta(TOTAL)
slide_objetivos(TOTAL)
slide_justificacion_alcance(TOTAL)
slide_estado_arte(TOTAL)
slide_marco_conductual(TOTAL)
slide_marco_vision_dl(TOTAL)
slide_metodologia(TOTAL)
slide_rf(TOTAL)
slide_rnf_reglas(TOTAL)
slide_riesgos_factibilidad(TOTAL)
slide_arquitectura(TOTAL)
slide_pipeline_ia(TOTAL)
slide_bd(TOTAL)
slide_vistas(TOTAL)
slide_metricas(TOTAL)
slide_conclusiones(TOTAL)
slide_bibliografia(TOTAL)
slide_cierre(TOTAL)

prs.save(OUT)
print("OK:", OUT)
