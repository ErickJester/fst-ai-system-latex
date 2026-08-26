#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generar_defensa.py
Agrega la sección 'Anexos (Material de Defensa)' a Presentacion_TT_B066_v17.pptx.
Fuente de datos: chapters/04_analisis.tex y chapters/05_diseno.tex (ya leídos y estructurados).
Salida: Presentacion_B066_Defensa.pptx
"""

# ── Dependencias ──────────────────────────────────────────────────────────────
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "lxml"])
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

# ── Rutas ─────────────────────────────────────────────────────────────────────
BASE   = r"C:\Users\ID634\Desktop\proyectos\latexFST\fst-ai-system-latex"
# NOTA: Presentacion_TT_B066_v17.pptx ya no vive en el repositorio ni en su
# historial. Para volver a ejecutar este script, copia el archivo desde el
# respaldo local (Escritorio/backup-tt-b066/presentaciones/) a docs/.
PPTX_IN  = os.path.join(BASE, "docs", "Presentacion_TT_B066_v17.pptx")
PPTX_OUT = os.path.join(BASE, "docs", "Presentacion_B066_Defensa.pptx")
DIAG     = os.path.join(BASE, "diagramas")
MERMAID  = os.path.join(BASE, "figures", "mermaid")

def img(*parts):
    """Devuelve la ruta de una imagen. Prefiere DIAG, cae a MERMAID."""
    p = os.path.join(DIAG, *parts)
    if os.path.exists(p):
        return p
    p2 = os.path.join(MERMAID, *parts)
    if os.path.exists(p2):
        return p2
    return None

# ── Colores del tema IPN-ESCOM (extraídos del PPTX original) ──────────────────
C_RED    = RGBColor(0x6B, 0x1F, 0x35)   # barra izquierda / títulos
C_GOLD   = RGBColor(0xC8, 0xA5, 0x35)   # línea / subrayados
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_PINK   = RGBColor(0xF7, 0xEE, 0xF0)   # filas pares
C_DARK   = RGBColor(0x1A, 0x20, 0x2C)   # texto cuerpo
C_LIGHT  = RGBColor(0xE8, 0xE8, 0xE8)   # bordes de tabla
C_LEVEL_HIGH = RGBColor(0xC0, 0x39, 0x2B)  # riesgo alto
C_LEVEL_MED  = RGBColor(0xD6, 0x89, 0x10)  # riesgo medio
C_LEVEL_LOW  = RGBColor(0x19, 0x6B, 0x24)  # riesgo bajo

# ── Geometría (EMU) ───────────────────────────────────────────────────────────
SLD_W   = 12192000
SLD_H   = 6858000
L_BAR   = 256032        # ancho de barra izquierda
CX      = L_BAR + 200000  # inicio del área de contenido
TITLE_Y = 480000
TITLE_H = 520000
LINE_Y  = 1060000
CON_Y   = 1150000       # inicio del contenido principal
FOOT_Y  = 6565392
FOOT_H  = 292608
CONTENT_W = SLD_W - CX - 100000

# ── Helpers de formato ────────────────────────────────────────────────────────
def _set_cell_text(cell, text, bold=False, sz=11, color=C_DARK,
                   align=PP_ALIGN.LEFT, font="Calibri"):
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    run.text = str(text)
    run.font.bold = bold
    run.font.size = Pt(sz)
    run.font.color.rgb = color
    run.font.name = font

def _style_cell(cell, fill_rgb=None, border_rgb=C_LIGHT, border_w=9525,
                text=None, bold=False, sz=9, color=C_DARK,
                align=PP_ALIGN.LEFT, font="Calibri"):
    """Aplica relleno y bordes a una celda con orden OOXML correcto (bordes antes de fill).
    Limpia elementos previos para evitar duplicados."""
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # 1. Limpiar fill y bordes previos (evita DUPLICATE y preserva orden)
    for tag in ('lnL', 'lnR', 'lnT', 'lnB', 'lnTlToBr', 'lnBlToTr',
                'solidFill', 'noFill', 'gradFill', 'blipFill', 'pattFill', 'grpFill'):
        for el in tcPr.findall(f'{{{NS}}}{tag}'):
            tcPr.remove(el)
    # 2. Bordes PRIMERO (orden correcto de esquema OOXML CT_TableCellProperties)
    if border_rgb is not None:
        for side in ('lnL', 'lnR', 'lnT', 'lnB'):
            ln = _etree.SubElement(tcPr, _qn(f'a:{side}'))
            ln.set('w', str(border_w))
            sf = _etree.SubElement(ln, _qn('a:solidFill'))
            sr = _etree.SubElement(sf, _qn('a:srgbClr'))
            sr.set('val', str(border_rgb).upper())
    # 3. Fill DESPUÉS de bordes
    if fill_rgb is not None:
        solidFill = _etree.SubElement(tcPr, _qn('a:solidFill'))
        srgb = _etree.SubElement(solidFill, _qn('a:srgbClr'))
        srgb.set('val', str(fill_rgb).upper())
    # 4. Texto opcional
    if text is not None:
        _set_cell_text(cell, text, bold=bold, sz=sz, color=color, align=align, font=font)

# ── Constructor de diapositivas ───────────────────────────────────────────────
class SlideBuilder:
    """Envuelve python-pptx para crear slides con el tema IPN-ESCOM."""

    def __init__(self, prs: Presentation):
        self.prs = prs
        self.blank = prs.slide_layouts[6]  # Blank layout

    # ── Slide separador (igual que "Conclusiones" del original) ───────────────
    def add_separator(self, title: str, subtitle: str = "Material de respaldo"):
        sld = self.prs.slides.add_slide(self.blank)
        sh = sld.shapes

        # BG blanco
        bg = sh.add_shape(1, Emu(0), Emu(0), Emu(SLD_W), Emu(SLD_H))
        bg.fill.solid(); bg.fill.fore_color.rgb = C_WHITE
        bg.line.fill.background()

        # Barra izquierda
        lb = sh.add_shape(1, Emu(0), Emu(0), Emu(L_BAR), Emu(SLD_H))
        lb.fill.solid(); lb.fill.fore_color.rgb = C_RED
        lb.line.fill.background()

        # Línea top/bot decorativas
        for y in (0, 6840000):
            bar = sh.add_shape(1, Emu(L_BAR), Emu(y), Emu(SLD_W - L_BAR), Emu(18000))
            bar.fill.solid(); bar.fill.fore_color.rgb = C_RED
            bar.line.fill.background()

        # Rombo dorado decorativo
        dmd = sh.add_shape(1, Emu(5745848), Emu(1700000), Emu(700000), Emu(700000))
        dmd.fill.solid()
        dmd.fill.fore_color.rgb = C_GOLD
        from pptx.util import Pt as _Pt
        from lxml import etree as _et
        solidFill = dmd._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                alpha = _et.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha.set('val', '40000')
        dmd.line.fill.background()

        # Título grande centrado
        tx = sh.add_textbox(Emu(600000), Emu(2200000), Emu(11000000), Emu(1600000))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = title
        run.font.bold = True; run.font.size = Pt(52)
        run.font.color.rgb = C_RED; run.font.name = "Georgia"

        # Subtítulo
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run(); run2.text = subtitle
        run2.font.size = Pt(28); run2.font.color.rgb = C_GOLD
        run2.font.name = "Georgia"

        # Footer
        ft = sh.add_shape(1, Emu(0), Emu(FOOT_Y), Emu(SLD_W), Emu(FOOT_H))
        ft.fill.solid(); ft.fill.fore_color.rgb = C_RED
        ft.line.fill.background()

        return sld

    # ── Slide con tabla ───────────────────────────────────────────────────────
    def add_table_slide(self, title: str, headers: list, rows: list,
                        col_widths=None, font_sz=9, header_sz=10,
                        slide_num=None):
        sld = self.prs.slides.add_slide(self.blank)
        sh = sld.shapes

        # Fondo + barra izquierda
        bg = sh.add_shape(1, Emu(0), Emu(0), Emu(SLD_W), Emu(SLD_H))
        bg.fill.solid(); bg.fill.fore_color.rgb = C_WHITE; bg.line.fill.background()

        lb = sh.add_shape(1, Emu(0), Emu(0), Emu(L_BAR), Emu(SLD_H))
        lb.fill.solid(); lb.fill.fore_color.rgb = C_RED; lb.line.fill.background()

        # Título
        tx = sh.add_textbox(Emu(CX), Emu(TITLE_Y), Emu(CONTENT_W), Emu(TITLE_H))
        tf = tx.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = title
        run.font.bold = True; run.font.size = Pt(22)
        run.font.color.rgb = C_RED; run.font.name = "Georgia"

        # Línea dorada
        ln = sh.add_shape(1, Emu(CX), Emu(LINE_Y), Emu(CONTENT_W - 200000), Emu(25000))
        ln.fill.solid(); ln.fill.fore_color.rgb = C_GOLD; ln.line.fill.background()

        # Tabla
        n_cols = len(headers)
        n_rows = len(rows) + 1  # +1 para encabezado
        tbl_w = SLD_W - CX - 100000
        tbl_h = SLD_H - CON_Y - FOOT_H - 100000

        if col_widths is None:
            col_widths = [tbl_w // n_cols] * n_cols
        else:
            # Normalizar a ancho total
            total = sum(col_widths)
            col_widths = [int(w / total * tbl_w) for w in col_widths]

        tbl = sh.add_table(n_rows, n_cols,
                            Emu(CX), Emu(CON_Y),
                            Emu(tbl_w), Emu(tbl_h)).table

        # Ajustar anchos de columna
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(w)

        # Encabezados
        for c, hdr in enumerate(headers):
            cell = tbl.cell(0, c)
            _style_cell(cell, fill_rgb=C_RED, text=hdr, bold=True, sz=header_sz,
                        color=C_WHITE, align=PP_ALIGN.CENTER)

        # Filas de datos
        for r, row in enumerate(rows):
            fill = C_PINK if r % 2 == 0 else C_WHITE
            for c, val in enumerate(row):
                cell = tbl.cell(r + 1, c)
                _style_cell(cell, fill_rgb=fill, text=val, sz=font_sz, align=PP_ALIGN.LEFT)

        # Footer + número
        ft = sh.add_shape(1, Emu(0), Emu(FOOT_Y), Emu(SLD_W), Emu(FOOT_H))
        ft.fill.solid(); ft.fill.fore_color.rgb = C_RED; ft.line.fill.background()

        if slide_num:
            num_tx = sh.add_textbox(Emu(10500000), Emu(30000), Emu(1550000), Emu(984885))
            num_tf = num_tx.text_frame
            np = num_tf.paragraphs[0]; np.alignment = PP_ALIGN.RIGHT
            nr = np.add_run(); nr.text = str(slide_num)
            nr.font.bold = True; nr.font.size = Pt(52)
            nr.font.color.rgb = C_RED; nr.font.name = "Georgia"

        return sld

    # ── Slide con imagen ──────────────────────────────────────────────────────
    def add_image_slide(self, title: str, img_path: str, caption: str = "",
                        slide_num=None):
        sld = self.prs.slides.add_slide(self.blank)
        sh = sld.shapes

        # Fondo + barra
        bg = sh.add_shape(1, Emu(0), Emu(0), Emu(SLD_W), Emu(SLD_H))
        bg.fill.solid(); bg.fill.fore_color.rgb = C_WHITE; bg.line.fill.background()
        lb = sh.add_shape(1, Emu(0), Emu(0), Emu(L_BAR), Emu(SLD_H))
        lb.fill.solid(); lb.fill.fore_color.rgb = C_RED; lb.line.fill.background()

        # Título
        tx = sh.add_textbox(Emu(CX), Emu(TITLE_Y), Emu(CONTENT_W), Emu(TITLE_H))
        tf = tx.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = title
        run.font.bold = True; run.font.size = Pt(22)
        run.font.color.rgb = C_RED; run.font.name = "Georgia"

        # Línea dorada
        ln = sh.add_shape(1, Emu(CX), Emu(LINE_Y), Emu(CONTENT_W - 200000), Emu(25000))
        ln.fill.solid(); ln.fill.fore_color.rgb = C_GOLD; ln.line.fill.background()

        # Imagen centrada
        img_y = CON_Y + 80000
        img_h = FOOT_Y - img_y - 200000
        img_w = SLD_W - CX - 200000

        if img_path and os.path.exists(img_path):
            from PIL import Image as PILImage
            try:
                with PILImage.open(img_path) as pil:
                    orig_w, orig_h = pil.size
                ratio = orig_w / orig_h
                # Mantener aspecto
                if img_w / img_h > ratio:
                    img_w_final = int(img_h * ratio)
                    img_h_final = img_h
                else:
                    img_w_final = img_w
                    img_h_final = int(img_w / ratio)
                left = CX + (img_w - img_w_final) // 2
                top  = img_y + (img_h - img_h_final) // 2
                sh.add_picture(img_path, Emu(left), Emu(top),
                               Emu(img_w_final), Emu(img_h_final))
            except Exception:
                # Fallback sin Pillow
                sh.add_picture(img_path, Emu(CX + 100000), Emu(img_y),
                               Emu(img_w - 200000), Emu(img_h))
        else:
            # Placeholder
            ph = sh.add_textbox(Emu(CX + 100000), Emu(img_y + 500000),
                                Emu(img_w - 200000), Emu(600000))
            ph_p = ph.text_frame.paragraphs[0]
            ph_p.alignment = PP_ALIGN.CENTER
            ph_r = ph_p.add_run()
            ph_r.text = f"[PENDIENTE: imagen no encontrada — {os.path.basename(img_path or 'N/A')}]"
            ph_r.font.size = Pt(12); ph_r.font.color.rgb = C_LEVEL_HIGH

        if caption:
            cap_tx = sh.add_textbox(Emu(CX), Emu(FOOT_Y - 250000),
                                    Emu(CONTENT_W), Emu(220000))
            cap_p = cap_tx.text_frame.paragraphs[0]
            cap_p.alignment = PP_ALIGN.CENTER
            cap_r = cap_p.add_run(); cap_r.text = caption
            cap_r.font.size = Pt(9); cap_r.font.color.rgb = C_DARK
            cap_r.font.italic = True

        # Footer
        ft = sh.add_shape(1, Emu(0), Emu(FOOT_Y), Emu(SLD_W), Emu(FOOT_H))
        ft.fill.solid(); ft.fill.fore_color.rgb = C_RED; ft.line.fill.background()

        if slide_num:
            num_tx = sh.add_textbox(Emu(10500000), Emu(30000), Emu(1550000), Emu(984885))
            num_tf = num_tx.text_frame
            np2 = num_tf.paragraphs[0]; np2.alignment = PP_ALIGN.RIGHT
            nr2 = np2.add_run(); nr2.text = str(slide_num)
            nr2.font.bold = True; nr2.font.size = Pt(52)
            nr2.font.color.rgb = C_RED; nr2.font.name = "Georgia"

        return sld


# ═════════════════════════════════════════════════════════════════════════════
# DATOS EXTRAÍDOS DE 04_analisis.tex y 05_diseno.tex  — TEXTO EXACTO DEL .tex
# ═════════════════════════════════════════════════════════════════════════════

# ── Requerimientos Funcionales (31 RF completos) ──────────────────────────────
RF_DATA = [
    # Autenticación
    ("RF-01", "Autenticación",    "Iniciar sesión con correo @ipn.mx y contraseña. El mensaje de error es genérico.", "Alta"),
    ("RF-02", "Autenticación",    "Cerrar sesión invalidando el token JWT desde cualquier pantalla.", "Alta"),
    ("RF-03", "Autenticación",    "Recuperar contraseña mediante enlace enviado al correo registrado.", "Media"),
    ("RF-04", "Autenticación",    "Dos roles: Investigador (consulta todos los experimentos) y Admin (más gestión de cuentas).", "Alta"),
    # Gestión de usuarios
    ("RF-05", "Gestión usuarios", "Admin crea, modifica y desactiva cuentas; usuario cambia contraseña temporal al primer ingreso.", "Alta"),
    ("RF-06", "Gestión usuarios", "El Investigador puede actualizar su nombre, correo y contraseña desde su perfil de usuario.", "Media"),
    ("RF-07", "Gestión usuarios", "No hay autoregistro público. Solo el Admin crea cuentas con correo institucional.", "Alta"),
    # Carga de video
    ("RF-08", "Carga de video",   "El sistema acepta únicamente archivos en formato .mp4. Cualquier otro se rechaza con mensaje claro.", "Alta"),
    ("RF-09", "Carga de video",   "Soporta hasta 2 videos por experimento: Día 1 (20 min) y Día 2 (5 min). Ambos son opcionales. El pipeline analiza el Día 2 completo; del Día 1 analiza unicamente los primeros 5 min y solo cuando corresponde al grupo control.", "Alta"),
    ("RF-10", "Carga de video",   "El sistema valida que el archivo sea un video reproducible antes de encolarlo. Si falla, lo rechaza con codigo de error y mensaje claro.", "Alta"),
    ("RF-11", "Carga de video",   "Al crear un experimento el investigador registra: nombre, fecha, numero de ratas (entre 1 y 4), tratamiento o condicion experimental, y notas adicionales opcionales.", "Alta"),
    ("RF-12", "Carga de video",   "El analisis se inicia automaticamente al completarse la carga y validacion del video. El investigador no puede iniciarlo, pausarlo, cancelarlo ni reiniciarlo desde la interfaz.", "Alta"),
    # Analisis
    ("RF-13", "Analisis IA",      "El pipeline ejecuta: preprocesamiento (CLAHE) → deteccion de cilindros (YOLOv8) → seguimiento (ByteTrack) → clasificacion de conductas (ResNet).", "Alta"),
    ("RF-14", "Analisis IA",      "Confianza minima de deteccion de cilindros: 0.70. Por debajo de ese umbral el pipeline se detiene (ver RN-11).", "Alta"),
    ("RF-15", "Analisis IA",      "Mientras el analisis corre, el investigador ve el progreso en tiempo real con barra de porcentaje y etapa activa. El analisis continua aunque se cierre la pestana.", "Alta"),
    ("RF-16", "Analisis IA",      "Si el analisis falla, el sistema muestra el codigo de error, la descripcion del problema y un boton para descargar el reporte de diagnostico en PDF.", "Alta"),
    ("RF-17", "Analisis IA",      "Clasifica cuadro a cuadro: nado activo, inmovilidad o escalamiento. Las conductas son mutuamente excluyentes. El buceo se clasifica como nado activo.", "Alta"),
    ("RF-18", "Analisis IA",      "Episodio valido >= 3 s consecutivos. Duracion menor se ignora (estandar del laboratorio).", "Alta"),
    # Resultados
    ("RF-19", "Resultados",       "Muestra tiempo total (s) y porcentaje de cada conducta por rata y sesion.", "Alta"),
    ("RF-20", "Resultados",       "El sistema ofrece un desglose de conductas por minuto para cada rata.", "Media"),
    ("RF-21", "Resultados",       "Cuando ambos videos del mismo experimento esten procesados, el sistema muestra una comparacion Dia 1 vs. Dia 2 con barras de distribucion por rata.", "Alta"),
    ("RF-22", "Resultados",       "Descarga CSV/XLSX con columnas Tiempo (min:seg) / Nado / Escalamiento / Inmovilidad. Reporte de diagnostico PDF por separado.", "Alta"),
    ("RF-31", "Resultados",       "Calcula estadisticas de grupo (media, desviacion estandar y varianza) por conducta entre ratas del mismo tratamiento.", "Media"),
    # Dashboard y retencion
    ("RF-24", "Dashboard",        "Muestra un aviso visible cuando un video esta a menos de 7 dias de ser eliminado automaticamente, indicando nombre del experimento y fecha de vencimiento exacta.", "Media"),
    ("RF-25", "Dashboard",        "El sistema borra automaticamente los archivos de video 30 dias despues de que el analisis termine exitosamente, notificando al investigador antes de que ocurra.", "Media"),
    ("RF-26", "Dashboard",        "Los resultados y reportes se conservan indefinidamente aunque el video original ya haya sido borrado.", "Alta"),
    ("RF-27", "Dashboard",        "El investigador puede eliminar un experimento desde el dashboard. La accion es permanente e irreversible y requiere confirmacion explicita (ver RN-09).", "Media"),
    # Administracion
    ("RF-28", "Administracion",   "El Admin puede ver la lista de todos los usuarios con correo, rol, estado (activo/inactivo), ultimo acceso y numero de experimentos. La lista es filtrable.", "Media"),
    ("RF-29", "Administracion",   "El Admin monitorea: uso de disco (%), espacio libre (GB), cola de analisis, analisis activos.", "Baja"),
    ("RF-30", "Administracion",   "El Admin puede ver y gestionar los experimentos de cualquier investigador (ver RN-08).", "Media"),
]

# ── Requerimientos No Funcionales (13 RNF con RNF-11) ─────────────────────────
RNF_DATA = [
    ("RNF-01", "Rendimiento",     "Respuesta en <= 3 s para el 95 % de peticiones de UI (dashboard, filtros, resultados)."),
    ("RNF-02", "Calidad IA",      "F1 >= 85 % por conducta sobre conjunto de prueba (nado, inmovilidad, escalamiento)."),
    ("RNF-03", "Disponibilidad",  "Disponible 24/7; >= 95 % de disponibilidad mensual. El analisis continua aunque el usuario cierre la pestana."),
    ("RNF-04", "Seguridad",       "Autenticacion JWT con tiempo de expiracion configurable. Solo usuarios autenticados acceden al sistema."),
    ("RNF-05", "Seguridad",       "Contrasenias cifradas con bcrypt+sal. Contrasenia temporal obliga cambio al primer ingreso."),
    ("RNF-06", "Usabilidad",      "Tasa de completacion de tarea >= 90 % sin asistencia tecnica en prueba con usuario real."),
    ("RNF-07", "Compatibilidad",  "Funciona en Chrome, Firefox y Edge (versiones actuales). Sin instalacion de software en el cliente."),
    ("RNF-08", "Almacenamiento",  "Videos se borran 30 dias tras analisis exitoso. Aviso con >= 7 dias de anticipacion."),
    ("RNF-09", "Almacenamiento",  "Alerta al Admin al 80 % de uso de disco. Limpieza automatica de videos mas antiguos al 90 %."),
    ("RNF-10", "Portabilidad",    "Despliegue con Docker Compose en servidor institucional provisto por la ESCOM sin dependencias adicionales."),
    ("RNF-11", "Escalabilidad",   "El modulo de clasificacion debe poder reconfigurarse para soportar paradigmas conductuales distintos al FST sin modificar la plataforma web. Criterio: diseno modular verificado en revision de arquitectura."),
    ("RNF-12", "Mantenibilidad",  "Codigo documentado por modulo (backend, pipeline IA, frontend). PEP 8 + ESLint."),
    ("RNF-13", "Concurrencia",    "Soporta hasta 4 usuarios activos simultaneamente. Analisis en cola secuencial (RN-04)."),
]

# ── Reglas de Negocio (texto exacto del .tex) ──────────────────────────────────
RN_DATA = [
    ("RN-01", "El mensaje de error de inicio de sesion es generico: no especifica si falla el correo o la contrasenia."),
    ("RN-02", "Un video pertenece a exactamente un experimento. No se puede reutilizar un video ya subido en un experimento distinto."),
    ("RN-03", "Un experimento tiene como maximo dos videos: uno del Dia 1 (20 min) y uno del Dia 2 (5 min)."),
    ("RN-04", "El analisis lo inicia el sistema automaticamente al terminar la carga. El investigador no puede iniciarlo, pausarlo, cancelarlo ni reiniciarlo desde la interfaz."),
    ("RN-05", "Videos se borran automaticamente 30 dias tras analisis exitoso. Irreversible. Aviso con >= 7 dias de anticipacion."),
    ("RN-06", "Resultados y reportes se conservan indefinidamente aunque el video original ya haya sido borrado."),
    ("RN-07", "En cada cuadro, a cada rata se le asigna exactamente una conducta. No coexisten dos conductas simultaneas."),
    ("RN-08", "Todos los investigadores activos ven todos los experimentos del laboratorio. El Admin ademas gestiona cuentas."),
    ("RN-09", "Borrar un experimento es permanente e irreversible: se eliminan el video, los resultados y los reportes. Requiere confirmacion explicita."),
    ("RN-10", "El sistema rechaza cualquier archivo que no sea .mp4 antes de intentar subirlo al servidor. La validacion de formato ocurre en el cliente."),
    ("RN-11", "Si la confianza de deteccion de los cuatro cilindros no alcanza 0.70, el pipeline se detiene y el sistema genera automaticamente un reporte de diagnostico descargable que indica el motivo del fallo. Los videos que superen ese umbral continuan el analisis de forma normal."),
    ("RN-12", "La comparacion Dia 1 vs. Dia 2 solo se genera cuando ambos videos del experimento estan procesados."),
    ("RN-13", "Si F1 < 85 % distinguiendo las 3 conductas, nado y escalamiento se reportan juntos como 'conducta activa'. El reporte indica el nivel aplicado."),
]

# ── Riesgos (R-01 a R-08 en orden del documento; R-06 aparece al final como Bajo) ─
RISKS_DATA = [
    ("R-01", "Alto",
     "Videos con reflejos, iluminacion variable o especimenes de pelaje claro que hacen fallar el pipeline (confianza < 0.70).",
     "Definir con el laboratorio los requisitos minimos de calidad de video. El sistema aplica CLAHE como primer intento; si la confianza sigue por debajo de 0.70, genera el reporte de diagnostico (RN-11)."),
    ("R-02", "Alto",
     "El dataset de cuadros etiquetados no alcanza para entrenar un clasificador con F1 >= 85 % objetivo.",
     "Ampliar el dataset anotando mas cuadros con BORIS y aplicar data augmentation. Fallback: reportar 'conducta activa' (RN-13)."),
    ("R-03", "Medio",
     "El laboratorio no proporciona los registros conductuales en el tiempo necesario para entrenamiento y validacion.",
     "Acordar un calendario de entrega con el laboratorio al inicio de la implementacion y aprovechar cada sesion con el Dr. Sandino."),
    ("R-04", "Medio",
     "El tiempo de analisis de un video de 5 min supera lo que el laboratorio puede asumir en su flujo de trabajo.",
     "Medir tiempos por etapa del pipeline. Evaluar reduccion de resolucion o uso de GPU. El umbral aceptable se define con el laboratorio."),
    ("R-05", "Medio",
     "La integracion del modulo de vision por computadora con el backend via cola asincrona resulta mas compleja de lo esperado.",
     "Definir contratos de API al inicio de la implementacion y ejecutar prueba end-to-end antes de agregar mas funcionalidad."),
    ("R-07", "Medio",
     "El laboratorio solicita cambios en los requerimientos a mitad del desarrollo.",
     "Diseno modular del clasificador para incorporar clases adicionales sin reescribir la arquitectura. Gestionar cambios dentro del proceso de revision de sprint conforme a Scrum."),
    ("R-08", "Medio",
     "El escalamiento y el nado activo se ven tan parecidos desde la vista lateral que el modelo no los distingue con suficiente precision.",
     "Usar en el entrenamiento unicamente clips inequivocos de escalamiento (movimiento de patas delanteras hacia las paredes). Fallback: reportar como 'conducta activa' (RN-13)."),
    ("R-06", "Bajo",
     "El disco del servidor se llena por acumulacion de videos durante el periodo de retencion de 30 dias.",
     "El sistema genera alertas al 80 % de uso de disco y ejecuta limpieza automatica al 90 % (RF-29, RNF-09)."),
]

# ── Costos ──────────────────────────────────────────────────────────────────────
COSTS_DATA = [
    ("FIJO",     "Recurso humano",          "2 desarrolladores x 500 h c/u @ $40/h",   "$40,000 MXN"),
    ("FIJO",     "GPU - Google Colab Pro",  "Plan anual, entrenamiento del clasificador", "$3,300 MXN"),
    ("FIJO",     "Conectividad (12 meses)", "Internet + electricidad durante desarrollo", "$4,500 MXN"),
    ("FIJO",     "Impresion y encuadernacion", "Dos ejemplares del documento recepcional", "$2,400 MXN"),
    ("FIJO",     "Tramites administrativos", "Titulacion y registro",                    "$800 MXN"),
    ("FIJO",     "Equipo y licencias",       "ENMyH-IPN provee equipo. Software open-source.", "$0"),
    ("VARIABLE", "Almacenamiento en nube",  "Si dataset supera capacidad gratuita",      "$0 - $1,500 MXN"),
    ("VARIABLE", "Insumos audiovisuales",   "Almacenamiento externo, adaptadores",       "$2,000 MXN"),
    ("VARIABLE", "Formacion tecnica",       "Cursos de vision por computadora / Docker",  "$2,000 MXN"),
    ("TOTAL",    "COSTO ESTIMADO TOTAL",    "Fijos $51,000 + Variables $5,500",          "~$56,500 MXN"),
]

# ── Diccionario de datos — 12 tablas, todas con id INTEGER (PK) autogenerado ───
DD_DATA = [
    # USUARIOS
    ("USUARIOS",         "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("USUARIOS",         "correo",                  "VARCHAR(120)",  "Correo institucional @ipn.mx (unico)"),
    ("USUARIOS",         "hash_contrasena",         "TEXT",          "Hash bcrypt+sal de la contrasenia"),
    ("USUARIOS",         "nombre",                  "VARCHAR(80)",   "Nombre completo del investigador"),
    ("USUARIOS",         "rol",                     "ENUM",          "investigador | admin"),
    ("USUARIOS",         "activo",                  "BOOLEAN",       "Cuenta habilitada / deshabilitada"),
    ("USUARIOS",         "debe_cambiar_contrasena", "BOOLEAN",       "True si la contrasenia es temporal"),
    ("USUARIOS",         "ultimo_acceso",           "TIMESTAMPTZ",   "Fecha y hora del ultimo ingreso"),
    # EXPERIMENTOS
    ("EXPERIMENTOS",     "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("EXPERIMENTOS",     "id_usuario",              "INTEGER (FK)",  "Referencia a USUARIOS.id"),
    ("EXPERIMENTOS",     "nombre",                  "VARCHAR(120)",  "Nombre descriptivo del experimento"),
    ("EXPERIMENTOS",     "fecha",                   "DATE",          "Fecha de realizacion"),
    ("EXPERIMENTOS",     "tratamiento",             "VARCHAR(120)",  "Molecula o condicion experimental"),
    ("EXPERIMENTOS",     "disposicion",             "INTEGER",       "Numero de ratas (1-4 segun diseno)"),
    ("EXPERIMENTOS",     "notas",                   "TEXT",          "Observaciones adicionales (opcional)"),
    # VIDEOS
    ("VIDEOS",           "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("VIDEOS",           "id_experimento",          "INTEGER (FK)",  "Referencia a EXPERIMENTOS.id"),
    ("VIDEOS",           "dia",                     "ENUM",          "dia1 | dia2"),
    ("VIDEOS",           "ruta_almacenamiento",     "TEXT",          "Ruta relativa en el volumen Docker"),
    ("VIDEOS",           "duracion_segundos",       "INTEGER",       "Duracion del video en segundos"),
    ("VIDEOS",           "fecha_eliminacion_prog.", "TIMESTAMPTZ",   "Fecha limite de retencion (30 dias)"),
    # SUJETOS
    ("SUJETOS",          "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("SUJETOS",          "id_experimento",          "INTEGER (FK)",  "Referencia a EXPERIMENTOS.id"),
    ("SUJETOS",          "indice_rata",             "SMALLINT",      "Numero del especimen dentro del experimento"),
    ("SUJETOS",          "etiqueta",                "VARCHAR(40)",   "Etiqueta descriptiva del especimen"),
    # ROIS
    ("ROIS",             "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("ROIS",             "id_sujeto",               "INTEGER (FK)",  "Referencia a SUJETOS.id"),
    ("ROIS",             "id_video",                "INTEGER (FK)",  "Referencia a VIDEOS.id"),
    ("ROIS",             "x",                       "INTEGER",       "Coordenada x del bounding box"),
    ("ROIS",             "y",                       "INTEGER",       "Coordenada y del bounding box"),
    ("ROIS",             "w",                       "INTEGER",       "Ancho del bounding box"),
    ("ROIS",             "h",                       "INTEGER",       "Alto del bounding box"),
    # CONFIGURACIONES_ANALISIS
    ("CONFIG_ANALISIS",  "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("CONFIG_ANALISIS",  "nombre_modelo",           "VARCHAR(80)",   "Nombre del modelo de clasificacion"),
    ("CONFIG_ANALISIS",  "hash_modelo",             "CHAR(64)",      "SHA-256 del archivo del modelo"),
    ("CONFIG_ANALISIS",  "version_pipeline",        "VARCHAR(20)",   "Version semantica del pipeline"),
    ("CONFIG_ANALISIS",  "umbral_confianza",        "NUMERIC(4,2)",  "Confianza minima de deteccion (0.70)"),
    ("CONFIG_ANALISIS",  "umbral_inmovil",          "NUMERIC(4,2)",  "Umbral de inmovilidad"),
    ("CONFIG_ANALISIS",  "umbral_desplazamiento",   "NUMERIC(6,2)",  "Umbral de desplazamiento"),
    ("CONFIG_ANALISIS",  "umbral_aspecto_trepada",  "NUMERIC(4,2)",  "Umbral de aspecto para escalamiento"),
    # TRABAJOS
    ("TRABAJOS",         "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("TRABAJOS",         "id_video",                "INTEGER (FK)",  "Referencia a VIDEOS.id"),
    ("TRABAJOS",         "id_config",               "INTEGER (FK)",  "Referencia a CONFIG_ANALISIS.id"),
    ("TRABAJOS",         "estado",                  "ENUM",          "pendiente|procesando|completado|error"),
    ("TRABAJOS",         "etapa_actual",            "VARCHAR(40)",   "Etapa activa del pipeline"),
    ("TRABAJOS",         "progreso_pct",            "SMALLINT",      "Progreso 0-100 (barra de estado)"),
    ("TRABAJOS",         "mensaje_error",           "TEXT",          "Descripcion del error si estado=error"),
    # ANIMALES (tabla de enlace TRABAJOS x SUJETOS)
    ("ANIMALES",         "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("ANIMALES",         "id_trabajo",              "INTEGER (FK)",  "Referencia a TRABAJOS.id"),
    ("ANIMALES",         "id_sujeto",               "INTEGER (FK)",  "Referencia a SUJETOS.id"),
    # RESULTADOS_COMPORTAMIENTO
    ("RES_COMP.",        "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("RES_COMP.",        "id_animal",               "INTEGER (FK)",  "Referencia a ANIMALES.id"),
    ("RES_COMP.",        "nado_s",                  "NUMERIC(6,2)",  "Segundos totales de nado activo"),
    ("RES_COMP.",        "inmovil_s",               "NUMERIC(6,2)",  "Segundos totales de inmovilidad"),
    ("RES_COMP.",        "escape_s",                "NUMERIC(6,2)",  "Segundos totales de escalamiento"),
    ("RES_COMP.",        "total_analizado_s",       "NUMERIC(6,2)",  "Duracion total analizada (s)"),
    # COMPORTAMIENTO_POR_MINUTO
    ("COMP_MIN.",        "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("COMP_MIN.",        "id_animal",               "INTEGER (FK)",  "Referencia a ANIMALES.id"),
    ("COMP_MIN.",        "minuto",                  "SMALLINT",      "Minuto analizado (1, 2, 3...)"),
    ("COMP_MIN.",        "nado_s",                  "NUMERIC(6,2)",  "Segundos de nado en ese minuto"),
    ("COMP_MIN.",        "inmovil_s",               "NUMERIC(6,2)",  "Segundos de inmovilidad en ese minuto"),
    ("COMP_MIN.",        "escape_s",                "NUMERIC(6,2)",  "Segundos de escalamiento en ese minuto"),
    # REPORTES
    ("REPORTES",         "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("REPORTES",         "id_trabajo",              "INTEGER (FK)",  "Referencia a TRABAJOS.id"),
    ("REPORTES",         "tipo_formato",            "ENUM",          "pdf | csv | xlsx"),
    ("REPORTES",         "ruta_archivo",            "TEXT",          "Ruta relativa en el volumen Docker"),
    ("REPORTES",         "fecha_creacion",          "TIMESTAMPTZ",   "Fecha y hora de generacion"),
    # NOTIFICACIONES
    ("NOTIFICACIONES",   "id",                      "INTEGER (PK)",  "Identificador unico autogenerado"),
    ("NOTIFICACIONES",   "id_usuario",              "INTEGER (FK)",  "Referencia a USUARIOS.id"),
    ("NOTIFICACIONES",   "id_experimento",          "INTEGER (FK)",  "Referencia a EXPERIMENTOS.id (opcional)"),
    ("NOTIFICACIONES",   "tipo",                    "VARCHAR(40)",   "analisis_ok|error|disco|expiracion"),
    ("NOTIFICACIONES",   "mensaje",                 "TEXT",          "Texto del aviso al investigador"),
    ("NOTIFICACIONES",   "leida",                   "BOOLEAN",       "True si el usuario ya vio el aviso"),
    ("NOTIFICACIONES",   "creada_en",               "TIMESTAMPTZ",   "Fecha y hora de creacion"),
]

# ── Stack tecnológico (sin versiones no documentadas en el .tex) ──────────────
STACK_DATA = [
    ("Sistema / Entorno",     "Servidor institucional ESCOM (Docker instalado)", "El laboratorio no instala nada mas alla de Docker"),
    ("Contenedores",          "Docker + Docker Compose",  "Orquestacion de 4 servicios: frontend, backend, worker, BD"),
    ("Lenguaje backend",      "Python",                   "API REST Flask y worker de analisis asincrono"),
    ("Framework web",         "Flask",                    "API REST; autenticacion JWT con Flask-JWT-Extended"),
    ("ORM",                   "SQLAlchemy",               "Mapeo objeto-relacional para PostgreSQL"),
    ("Base de datos",         "PostgreSQL",               "BD relacional con transacciones ACID, codigo abierto"),
    ("Lenguaje frontend",     "JavaScript / TypeScript",  "React.js + Vite; SPA con barra de progreso en tiempo real"),
    ("Deteccion objetos",     "YOLOv8 (Ultralytics)",    "Localizacion de cilindros y especimenes en el video"),
    ("Seguimiento",           "ByteTrack",                "Asociacion temporal de detecciones entre cuadros"),
    ("Clasificacion",         "ResNet-18 / ResNet-50",    "Clasificador de conductas (pre-entrenado en ImageNet)"),
    ("Framework DL",          "PyTorch",                  "Entrenamiento e inferencia del clasificador"),
    ("Vision por computadora","OpenCV 4.x",               "Decodificacion de video, CLAHE, filtros morfologicos"),
    ("Computo GPU",           "Google Colab Pro",         "Entrenamiento del clasificador en GPU (fase TT-II)"),
    ("Anotacion",             "BORIS",                    "Behavioral Observation Research Interactive Software"),
    ("Calidad de codigo",     "PEP 8 + ESLint",           "Revision automatica en Python y JavaScript"),
]

# ── Trazabilidad RF → Módulo ──────────────────────────────────────────────────
TRACE_DATA = [
    ("RF-01 / RF-02 / RF-03", "Autenticacion",       "Backend (JWT)",   "USUARIOS",               "—"),
    ("RF-04 / RF-05 / RF-06 / RF-07", "Gestion usuarios", "Backend (roles)", "USUARIOS",           "—"),
    ("RF-08 / RF-09 / RF-10 / RF-11", "Carga de video", "Backend",        "VIDEOS / EXPERIMENTOS", "—"),
    ("RF-12",                 "Encolamiento",        "Backend + Worker", "TRABAJOS",               "—"),
    ("RF-13 / RF-14",         "Deteccion cilindros", "Worker (YOLO)",   "TRABAJOS / ROIS",         "YOLOv8"),
    ("RF-15 / RF-16",         "Progreso y errores",  "Backend + WS",    "TRABAJOS",               "—"),
    ("RF-17 / RF-18",         "Clasificacion cond.", "Worker (ResNet)", "RES_COMP. / ANIMALES",    "ResNet-18/50"),
    ("RF-19 / RF-20 / RF-21", "Visualizacion resul.","Backend",         "RES_COMP. / COMP_MIN.",   "—"),
    ("RF-22 / RF-31",         "Exportacion datos",   "Backend",         "REPORTES",               "pandas / openpyxl"),
    ("RF-24 / RF-25 / RF-26 / RF-27", "Dashboard / retencion","Backend (cron)", "EXPERIMENTOS / VIDEOS", "—"),
    ("RF-28 / RF-29 / RF-30", "Panel administracion","Backend",         "NOTIFICACIONES",          "—"),
]


# ═════════════════════════════════════════════════════════════════════════════
# MAIN — Construcción de la presentación
# ═════════════════════════════════════════════════════════════════════════════
def main():
    print(f"Abriendo: {PPTX_IN}")
    prs = Presentation(PPTX_IN)
    sb  = SlideBuilder(prs)

    total_existente = len(prs.slides)
    n = total_existente  # contador de número de diapositiva

    print(f"Diapositivas existentes: {total_existente}")
    print("Agregando sección de Anexos...")

    # ── 1. SEPARADOR ──────────────────────────────────────────────────────────
    sb.add_separator("Anexos", "Material de Defensa")
    n += 1; print(f"  [{n}] Separador: Anexos")

    # ── 2. DICCIONARIO DE DATOS (3 slides — 12 tablas) ────────────────────────
    t3 = len(DD_DATA) // 3
    dd_parts = [DD_DATA[:t3], DD_DATA[t3:2*t3], DD_DATA[2*t3:]]
    for part, rows in enumerate(dd_parts, 1):
        sb.add_table_slide(
            f"Diccionario de Datos ({part}/3)",
            ["Tabla", "Atributo", "Tipo de dato", "Descripcion"],
            rows,
            col_widths=[2.2, 2.8, 2, 3.5],
            font_sz=8,
            slide_num=n + 1
        )
        n += 1; print(f"  [{n}] Diccionario de Datos parte {part}")

    # ── 3. CONFIGURACIÓN DE ENTORNO ──────────────────────────────────────────
    sb.add_table_slide(
        "Configuración de Entorno y Stack Tecnológico",
        ["Componente", "Tecnología / Versión", "Rol en el sistema"],
        STACK_DATA,
        col_widths=[3, 3, 5],
        font_sz=9,
        slide_num=n + 1
    )
    n += 1; print(f"  [{n}] Stack tecnológico")

    # ── 4. TRAZABILIDAD RF → MÓDULOS ─────────────────────────────────────────
    sb.add_table_slide(
        "Trazabilidad: Requerimientos → Módulos",
        ["RF(s)", "Módulo funcional", "Capa técnica", "Tablas BD", "Librería/Tech"],
        TRACE_DATA,
        col_widths=[3, 2.5, 2, 2.5, 2],
        font_sz=8,
        slide_num=n + 1
    )
    n += 1; print(f"  [{n}] Trazabilidad")

    # ── 5. MATRIZ DE RIESGOS ─────────────────────────────────────────────────
    risk_rows = [(r[0], r[1], r[2], r[3]) for r in RISKS_DATA]
    sld_risk = sb.add_table_slide(
        "Matriz de Riesgos Críticos",
        ["ID", "Nivel", "Descripción", "Plan de mitigación"],
        risk_rows,
        col_widths=[1, 1.2, 4.5, 5.5],
        font_sz=9,
        slide_num=n + 1
    )
    # Colorear celda "Nivel" según valor
    tbl_risk = None
    for shape in sld_risk.shapes:
        if shape.shape_type == 19:
            tbl_risk = shape.table; break
    if tbl_risk:
        LEVEL_COLORS = {"Alto": C_LEVEL_HIGH, "Medio": C_LEVEL_MED, "Bajo": C_LEVEL_LOW}
        for r, row in enumerate(RISKS_DATA):
            nivel = row[1]
            cell = tbl_risk.cell(r + 1, 1)
            _style_cell(cell, fill_rgb=LEVEL_COLORS.get(nivel, C_DARK),
                        text=nivel, bold=True, sz=9, color=C_WHITE, align=PP_ALIGN.CENTER)
    n += 1; print(f"  [{n}] Matriz de riesgos")

    # ── 6. REQUERIMIENTOS FUNCIONALES (3 slides — 30 RF) ─────────────────────
    rf_t3 = len(RF_DATA) // 3
    rf_parts = [RF_DATA[:rf_t3], RF_DATA[rf_t3:2*rf_t3], RF_DATA[2*rf_t3:]]
    for part, rows in enumerate(rf_parts, 1):
        formatted = [(r[0], r[1], r[2], r[3]) for r in rows]
        sb.add_table_slide(
            f"Requerimientos Funcionales ({part}/3)",
            ["ID", "Modulo", "Descripcion", "Prio."],
            formatted,
            col_widths=[1.0, 1.8, 8.0, 1.2],
            font_sz=7,
            slide_num=n + 1
        )
        n += 1; print(f"  [{n}] RF parte {part}")

    # ── 7. REQUERIMIENTOS NO FUNCIONALES ────────────────────────────────────
    rnf_rows = [(r[0], r[1], r[2]) for r in RNF_DATA]
    sb.add_table_slide(
        "Requerimientos No Funcionales",
        ["ID", "Categoría", "Descripción y métrica"],
        rnf_rows,
        col_widths=[1.2, 2, 8.8],
        font_sz=8,
        slide_num=n + 1
    )
    n += 1; print(f"  [{n}] RNF")

    # ── 8. REGLAS DE NEGOCIO ─────────────────────────────────────────────────
    sb.add_table_slide(
        "Reglas de Negocio",
        ["ID", "Regla"],
        RN_DATA,
        col_widths=[1.2, 11],
        font_sz=8,
        slide_num=n + 1
    )
    n += 1; print(f"  [{n}] Reglas de negocio")

    # ── 9. ESTIMACIÓN DE COSTOS ──────────────────────────────────────────────
    sld_cost = sb.add_table_slide(
        "Estimación de Costos del Proyecto",
        ["Tipo", "Rubro", "Detalle", "Monto (MXN)"],
        COSTS_DATA,
        col_widths=[1.5, 3, 5, 2.5],
        font_sz=9,
        slide_num=n + 1
    )
    # Resaltar fila TOTAL
    for shape in sld_cost.shapes:
        if shape.shape_type == 19:
            tbl_cost = shape.table
            for r in range(tbl_cost.rows.__len__()):
                try:
                    if tbl_cost.cell(r, 0).text_frame.text == "TOTAL":
                        for c in range(4):
                            prev_text = tbl_cost.cell(r, c).text_frame.text
                            _style_cell(tbl_cost.cell(r, c), fill_rgb=C_RED,
                                        text=prev_text, bold=True, sz=10, color=C_WHITE)
                except Exception:
                    pass
            break
    n += 1; print(f"  [{n}] Costos")

    # ── 10+. DIAGRAMAS (imágenes) — todos los del capítulo de diseño ────────
    diag_slides = [
        # ── Casos de Uso ──────────────────────────────────────────────────────
        ("Casos de Uso — Visión General",
         img("cu_vision_general.png"),
         "Actores: Investigador y Administrador · Relaciones de alto nivel"),
        ("Casos de Uso — Paquete 1: Autenticación y Gestión de Usuarios",
         img("cu_paquete1.png"),
         "Registro, inicio de sesión, cambio de contraseña, edición de perfil"),
        ("Casos de Uso — Paquete 2: Gestión de Experimentos y Carga de Video",
         img("cu_paquete2.png"),
         "Crear experimento, cargar video .mp4, validación y encolamiento"),
        ("Casos de Uso — Paquete 3: Análisis Conductual",
         img("cu_paquete3.png"),
         "Ejecución del pipeline YOLOv8 + ByteTrack + ResNet en el Worker"),
        ("Casos de Uso — Paquete 4: Resultados y Reportes",
         img("cu_paquete4.png"),
         "Consulta de resultados, exportación CSV / XLSX / PDF"),
        ("Casos de Uso — Paquete 5: Dashboard, Notificaciones y Administración",
         img("cu_paquete5.png"),
         "Monitoreo del sistema, gestión de cuentas, alertas de almacenamiento"),

        # ── Secuencias: Autenticación y Usuarios ─────────────────────────────
        ("Secuencia — Registro de Investigador",
         img("seq_registro.png"),
         "El Administrador crea la cuenta; el usuario activa con contraseña temporal"),
        ("Secuencia — Inicio de Sesión (JWT)",
         img("seq_login.png"),
         "Validación de credenciales, generación y almacenamiento de token JWT"),
        ("Secuencia — Cierre de Sesión",
         img("seq_logout.png"),
         "Invalidación del token en el cliente; redirección al login"),
        ("Secuencia — Cambio de Contraseña",
         img("seq_cambio_pass_inv.png"),
         "Verificación de contraseña actual + confirmación de la nueva"),
        ("Secuencia — Gestión de Usuarios (Administrador)",
         img("seq_gestion_usuarios.png"),
         "Aprobar solicitudes, editar cuentas, desactivar accesos"),
        ("Secuencia — Edición de Perfil",
         img("seq_perfil.png"),
         "Actualización de nombre, correo o contraseña por cualquier usuario activo"),
        ("Secuencia — Notificaciones del Sistema",
         img("seq_notificaciones.png"),
         "Avisos de almacenamiento y expiración de videos (RN-11, RNF-10)"),

        # ── Secuencias: Pipeline (separadas en carga + análisis) ─────────────
        ("Secuencia — Carga de Video",
         img("seq_carga.png"),
         "Investigador sube .mp4 → validación → almacenamiento → encolamiento"),
        ("Secuencia — Análisis Automático (Worker)",
         img("seq_analisis.png"),
         "Worker toma tarea de la cola → YOLOv8 → ByteTrack → ResNet → persiste resultados"),
        ("Secuencia — Seguimiento de Progreso",
         img("seq_progreso.png"),
         "Polling del frontend al backend; estados: en cola / procesando / completado"),
        ("Secuencia — Manejo de Error en Análisis",
         img("seq_error.png"),
         "Camino alternativo: confianza < 0.70 o fallo del Worker → notificación"),
        ("Secuencia — Consulta de Resultados",
         img("seq_resultados.png"),
         "Tiempos por conducta y espécimen; estadísticos de grupo (RF-31)"),
        ("Secuencia — Descarga de Reportes",
         img("seq_reportes.png"),
         "Exportación CSV / XLSX / PDF a demanda del investigador"),

        # ── Arquitectura y Pipeline ───────────────────────────────────────────
        ("Arquitectura de Software (Capas)",
         img("arquitectura_software.png"),
         "Frontend React · Backend Flask · Worker Python · PostgreSQL · Volumen Docker"),
        ("Arquitectura del Sistema (Despliegue con Docker Compose)",
         img("arquitectura.png"),
         "4 contenedores: frontend :3000, backend :5000, worker, db :5432"),
        ("Pipeline de Análisis Conductual",
         img("pipeline.png"),
         "Preprocesamiento CLAHE → YOLOv8 (detección) → ByteTrack (seguimiento) → ResNet-50 (clasificación)"),

        # ── Modelos de Datos ──────────────────────────────────────────────────
        ("Modelo Entidad-Relación (Conceptual)",
         img("entidadRelacion.png"),
         "Entidades: USUARIOS, EXPERIMENTOS, VIDEOS, TRABAJOS, RESULTADOS"),
        ("Modelo Relacional (Esquema Físico)",
         img("er2.png"),
         "12 tablas con FKs, tipos de dato PostgreSQL y restricciones de integridad"),
        ("Diagrama de Clases del Sistema",
         img("clases.png"),
         "Clases centrales: Experimento, Video, Trabajo, ResultadoComportamiento"),
    ]

    for title, img_path, caption in diag_slides:
        sb.add_image_slide(title, img_path, caption, slide_num=n + 1)
        status = "OK" if (img_path and os.path.exists(img_path)) else "IMAGEN NO ENCONTRADA"
        n += 1; print(f"  [{n}] {title} — {status}")

    # ── Guardar ───────────────────────────────────────────────────────────────
    prs.save(PPTX_OUT)
    print(f"\nGuardado: {PPTX_OUT}")
    print(f"  Total diapositivas: {len(prs.slides)}")
    print(f"  Nuevas diapositivas agregadas: {len(prs.slides) - total_existente}")


if __name__ == "__main__":
    main()
